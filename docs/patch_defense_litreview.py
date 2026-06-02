"""
Populate slides 7 (Related Works) and 8 (Literature Review) of
HVAC_DRL_MORL_defense_full.pptx with the structured citation summary
provided by the author.

The two slides were previously left empty by design (the author wanted
to author the literature content themselves). They are now filled with
the author-supplied Russian-language summary, organised into five
research-program threads:

  Slide 7 (Related Works):
    1. Theoretical foundation and research gaps (Articles 1, 1.2, 33)
    2. Surrogate methodology, Block 1            (Hou-Evins, 14, 7)
    3. Control algorithms and safety, Block 2    (Articles 7, 15, 27, 28)

  Slide 8 (Literature Review):
    4. Feature design and MORL                    (Articles 22, 24, Max-Min, FastML)
    5. Transferability and future directions     (Articles 11, 30, 35)
    Synthesis  — what this work uniquely contributes

The placeholder shapes (dashed box and grey "[to be added]" text) are
removed; the title bar, subtitle, and slide footer are preserved.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
IN_PATH = ROOT / "docs" / "HVAC_DRL_MORL_defense_full.pptx"
OUT_PATH = IN_PATH

INK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xF4, 0xA2, 0x61)
BLUE = RGBColor(0x1E, 0x60, 0x91)
SUCCESS = RGBColor(0x0F, 0x8A, 0x5F)
PURPLE = RGBColor(0x5A, 0x4E, 0x7C)
WARN = RGBColor(0xC1, 0x12, 0x1F)
MUTED = RGBColor(0x6B, 0x70, 0x80)
LIGHT_BG = RGBColor(0xF7, 0xF6, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def remove_body_shapes(slide):
    """Aggressive cleanup: remove ALL shapes in the body area (top >= 1.3")
    except the subtitle (light italic, top ~ 0.9"-1.3") and the bottom
    footer pagination textbox. Preserves: title bar, title text, divider
    line, subtitle, footer text 'HVAC DRL/MORL Defense ... Slide N'.
    Idempotent: safe to run on already-patched slides."""
    to_remove = []
    for sh in slide.shapes:
        try:
            top_inches = sh.top / 914400
            height_inches = sh.height / 914400
        except Exception:
            continue
        # Title bar / divider / title text live in top 1.3"
        if top_inches < 1.3:
            continue
        # Bottom footer (page number etc.) lives below 6.95"
        if top_inches > 6.95:
            continue
        # Preserve text that is the slide footer
        if sh.has_text_frame:
            t = (sh.text_frame.text or "").strip()
            if "HVAC DRL/MORL Defense" in t:
                continue
        to_remove.append(sh)
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)


def update_subtitle(slide, new_text):
    """Find the subtitle textbox (small italic text in the top 1.3" of the
    slide that isn't the title) and replace its content. Idempotent."""
    title_text_starts = ("Related Works", "Literature Review")
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        try:
            top_inches = sh.top / 914400
        except Exception:
            continue
        if top_inches >= 1.3:
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        # skip the big title text
        if any(t.startswith(prefix) for prefix in title_text_starts) and len(t) < 30:
            continue
        # this is the subtitle — replace it
        tf = sh.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = new_text
        r.font.name = "Calibri"
        r.font.size = Pt(13)
        r.font.italic = True
        r.font.color.rgb = MUTED
        return


def add_text(slide, text, *, left, top, width, height,
             font_size=11, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_section_header(slide, label, *, left, top, width, color=BLUE):
    """Coloured header strip for a literature-review section."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(0.08), Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, label, left=left + 0.14, top=top - 0.02,
             width=width - 0.14, height=0.36,
             font_size=12, bold=True, color=color)


def add_paragraph_list(slide, items, *, left, top, width, height,
                       font_size=10, color=INK):
    """Add a textbox where each item is a paragraph with citation tag (bold)
    and explanation (regular)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, item in enumerate(items):
        # item is (citation, body) tuple
        cite, body = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        # bullet
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = "Calibri"
        r0.font.size = Pt(font_size)
        r0.font.bold = True
        r0.font.color.rgb = color
        # citation
        r1 = p.add_run()
        r1.text = cite + "  "
        r1.font.name = "Calibri"
        r1.font.size = Pt(font_size)
        r1.font.bold = True
        r1.font.color.rgb = color
        # body
        r2 = p.add_run()
        r2.text = body
        r2.font.name = "Calibri"
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color


# ---------------------------------------------------------------------------
# Slide 7 — Related Works
# ---------------------------------------------------------------------------

def patch_slide_7(prs: Presentation):
    s = prs.slides[6]  # 0-indexed
    remove_body_shapes(s)
    update_subtitle(s,
        "Theoretical foundation, surrogate methodology, and control algorithms")

    # Section 1
    add_section_header(s, "1.  Theoretical foundation and research gaps",
                       left=0.4, top=1.40, width=12.5, color=BLUE)
    add_paragraph_list(s,
        items=[
            ("Article 1 (Al Sayed et al., 2024):",
             "classifies building systems as POMDPs; highlights the deployment "
             "gap — only 23% of RL studies reach real buildings due to safety risks."),
            ("Article 1.2 (Xu & Ghahramani, 2026):",
             "stresses the exploration burden imposed on occupants and the need "
             "for RL adaptivity to personal comfort preferences."),
            ("Article 33 (Survey):",
             "in-depth analysis of distributional shift — how surrogate errors "
             "lead to policy-transfer failures in real-world deployment."),
        ],
        left=0.55, top=1.78, width=12.4, height=1.45,
        font_size=10)

    # Section 2
    add_section_header(s, "2.  Surrogate development methodology (Block 1)",
                       left=0.4, top=3.30, width=12.5, color=SUCCESS)
    add_paragraph_list(s,
        items=[
            ("Hou & Evins (2024):",
             "the project's core methodological standard. The entire surrogate "
             "development pipeline follows the 4-stage protocol with numerical "
             "evidence (Reporting Level 3) — supplementary tables S1-S11."),
            ("Article 14 (Coraci et al., 2025) + Article 7 (Wang et al., 2025):",
             "justify the use of first-order RC models and staged calibration. "
             "These ideas underpin our Stage A/B/C pipeline that identified C_zon."),
        ],
        left=0.55, top=3.68, width=12.4, height=1.10,
        font_size=10)

    # Section 3
    add_section_header(s, "3.  Control algorithms and safety (Block 2)",
                       left=0.4, top=4.85, width=12.5, color=ACCENT)
    add_paragraph_list(s,
        items=[
            ("Article 7 (Wang et al., 2025):",
             "introduced the integral safety metric m_s, accounting for both "
             "duration and severity of violations. Adopted as the final quality "
             "criterion across all our controllers."),
            ("Article 15 (Hedayat et al., 2025):",
             "set the gold standard for statistical rigour (50 seeds, "
             "confidence intervals) and sanity-check protocols via PID baselines."),
            ("Article 27 (Liao et al., 2025):",
             "basis of our hierarchical RL (HDRL) implementation — separating "
             "mode selection from setpoint tracking."),
            ("Article 28 (Savino et al., 2025):",
             "reference ASHRAE G36 control sequences for benchmarking "
             "low-level controllers."),
        ],
        left=0.55, top=5.23, width=12.4, height=1.60,
        font_size=10)


# ---------------------------------------------------------------------------
# Slide 8 — Literature Review
# ---------------------------------------------------------------------------

def patch_slide_8(prs: Presentation):
    s = prs.slides[7]
    remove_body_shapes(s)
    update_subtitle(s,
        "Feature design, MORL, transferability, and the synthesis "
        "that distinguishes this work")

    # Section 4
    add_section_header(s, "4.  Feature design and multi-objective optimisation (MORL)",
                       left=0.4, top=1.40, width=12.5, color=PURPLE)
    add_paragraph_list(s,
        items=[
            ("Article 22 (Gao et al., 2024) + Article 24 (Sun et al., 2024):",
             "demonstrated the critical role of predictive information (weather "
             "forecasts) and state history in converting a POMDP into an MDP. "
             "Motivated our 17-D interface — comfort-violation rate reduced "
             "from 74.5% to 4.9%."),
            ("Multi-Objective RL with Max-Min Criterion:",
             "theoretical foundation for reward scalarisation in the search "
             "for Pareto trade-offs between energy and comfort."),
            ("FastML:",
             "inspired the use of radar charts for visualising KPI balance "
             "(see slide 32 — MORL 5-D vs 17-D radar)."),
        ],
        left=0.55, top=1.78, width=12.4, height=1.45,
        font_size=10)

    # Section 5
    add_section_header(s, "5.  Transferability and future directions (Block 3 and beyond)",
                       left=0.4, top=3.30, width=12.5, color=BLUE)
    add_paragraph_list(s,
        items=[
            ("Article 11 (Hou et al., 2024):",
             "foundational work on multi-source transfer (MTL-DRL). Forms the "
             "basis of the Block 3 plan — training across families of test "
             "cases with explicit recalibration regimes."),
            ("Article 30 (Dreamer / Hafner et al., 2025):",
             "motivates the Dyna-style architecture in which the agent learns "
             "inside the latent imagination of a world model. Our hybrid "
             "backend (v3 + v3.5 disagreement) is a direct prototype of such a system."),
            ("Article 35 (LEGION / Bing et al., 2025):",
             "proposes lifelong-learning methods to prevent catastrophic "
             "forgetting when controlling many heterogeneous buildings."),
        ],
        left=0.55, top=3.68, width=12.4, height=1.45,
        font_size=10)

    # Synthesis box
    syn_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.4), Inches(5.25),
                                 Inches(12.5), Inches(1.45))
    syn_box.fill.solid()
    syn_box.fill.fore_color.rgb = LIGHT_BG
    syn_box.line.color.rgb = SUCCESS
    syn_box.line.width = Pt(1.5)

    add_text(s, "Synthesis — what this work uniquely contributes",
             left=0.55, top=5.32, width=12.2, height=0.35,
             font_size=12, bold=True, color=SUCCESS)
    add_text(s,
        "This work synthesises three threads: statistical rigour "
        "(Article 15 — 50-seed standard; we use 5 seeds plus a replay test), "
        "physical accuracy (Articles 7 and 14 — RC modelling; we apply staged "
        "inverse calibration with an identifiable C_zon), and protocol "
        "discipline (Hou & Evins 2024 — we implement all four stages at "
        "Reporting Level 3). This synthesis, together with four pre-registered "
        "audit anchors in git, makes the work methodologically unique in the "
        "current DRL-HVAC research landscape.",
        left=0.55, top=5.68, width=12.2, height=0.95,
        font_size=10, italic=False, color=INK)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"[INFO] reading  {IN_PATH}")
    prs = Presentation(str(IN_PATH))
    print(f"[INFO] slides   {len(prs.slides)}")
    patch_slide_7(prs)
    print(f"[OK]   patched slide 7 — Related Works (sections 1, 2, 3)")
    patch_slide_8(prs)
    print(f"[OK]   patched slide 8 — Literature Review (sections 4, 5 + synthesis)")
    prs.save(str(OUT_PATH))
    print(f"[OK]   wrote    {OUT_PATH}")


if __name__ == "__main__":
    main()
