"""
Patch the formula slides of HVAC_DRL_MORL_defense_full.pptx into a more
readable LaTeX-style mathematical typography.

PowerPoint does not natively render LaTeX, but we can substantially improve
readability with:
  - Unicode subscripts/superscripts where digits are 0-9 (a0 -> a₀, ^2 -> ²)
  - Greek letters and operators (already Unicode-safe)
  - Italic typeface for variables, regular for operators (run-level)
  - Multi-line equation layout with consistent indentation
  - Monospace alignment via Cambria Math / Consolas font where helpful

Touched slides:
  11  Hybrid loss              — multi-line ℒ_total with clear normed terms
  13  Surrogate equations      — v3, v3.5, hybrid loss in three blocks
  14  Control interface, reward, and safety metric — action / observation /
                                  reward / m_s with proper math typography

Slide numbering follows the existing 53-slide deck.
"""

from __future__ import annotations

from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
IN_PATH = ROOT / "docs" / "HVAC_DRL_MORL_defense_full.pptx"
OUT_PATH = IN_PATH                                  # overwrite in place

INK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xF4, 0xA2, 0x61)
BLUE = RGBColor(0x1E, 0x60, 0x91)
SUCCESS = RGBColor(0x0F, 0x8A, 0x5F)
MUTED = RGBColor(0x6B, 0x70, 0x80)

EQUATION_FONT = "Cambria Math"
PROSE_FONT = "Calibri"

# Math symbols used throughout
SUB = str.maketrans("0123456789", "₀₁₂₃₄₅₆₇₈₉")
SUP = str.maketrans("0123456789+-=()", "⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾")


def sub(s: str) -> str:
    return s.translate(SUB)


def sup(s: str) -> str:
    return s.translate(SUP)


# ---------------------------------------------------------------------------
# Helpers for cleaning a slide
# ---------------------------------------------------------------------------

def remove_shape_by_name_substring(slide, substrings):
    """Remove any shape whose internal name contains any of the substrings."""
    to_remove = []
    for sh in slide.shapes:
        if any(needle.lower() in (sh.name or "").lower() for needle in substrings):
            to_remove.append(sh)
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)


def remove_tables(slide):
    """Remove any TABLE shapes from a slide."""
    to_remove = []
    for sh in slide.shapes:
        if sh.has_table:
            to_remove.append(sh)
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)


def remove_textbox_by_text_prefix(slide, prefixes):
    """Remove TEXT_BOX shapes whose text starts with any of the prefixes."""
    to_remove = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if any(t.startswith(p) for p in prefixes):
                to_remove.append(sh)
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)


def add_equation_textbox(slide, *, left, top, width, height, lines,
                         font=EQUATION_FONT, base_size=15,
                         color=INK, align=PP_ALIGN.LEFT):
    """Add a textbox with multi-line equations.

    `lines` is a list of either:
      - a plain string  -> rendered as a single run
      - a list of (segment, style) tuples where style is one of:
          'eq'    : equation body text (math font, color INK)
          'var'   : italic variable
          'op'    : regular operator (math font)
          'label' : muted regular sans-serif (for 'where ...' notes)
          'em'    : bold accent for headers within an equation
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)

        # Each line can be a string or a list of (text, style) tuples
        segments = line if isinstance(line, list) else [(line, "eq")]

        for seg_text, seg_style in segments:
            r = p.add_run()
            r.text = seg_text
            if seg_style == "eq":
                r.font.name = font
                r.font.size = Pt(base_size)
                r.font.color.rgb = color
            elif seg_style == "var":
                r.font.name = font
                r.font.size = Pt(base_size)
                r.font.italic = True
                r.font.color.rgb = color
            elif seg_style == "op":
                r.font.name = font
                r.font.size = Pt(base_size)
                r.font.color.rgb = color
            elif seg_style == "label":
                r.font.name = PROSE_FONT
                r.font.size = Pt(base_size - 2)
                r.font.italic = True
                r.font.color.rgb = MUTED
            elif seg_style == "em":
                r.font.name = PROSE_FONT
                r.font.size = Pt(base_size - 1)
                r.font.bold = True
                r.font.color.rgb = BLUE
            else:
                r.font.name = font
                r.font.size = Pt(base_size)
                r.font.color.rgb = color
    return tb


def add_section_header(slide, text, *, left, top, width, color=BLUE, size=15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = PROSE_FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return tb


# ---------------------------------------------------------------------------
# Slide 11 — Hybrid loss
# ---------------------------------------------------------------------------

def patch_slide_11(prs: Presentation):
    slide = prs.slides[10]  # zero-indexed; slide 11

    # Remove the old "L_total = ..." textbox + its wrapper rectangle
    # (we identify by the existing rounded rectangle that holds the equation
    # and the plain textbox that sits inside it)
    # Strategy: drop shapes by name `Rectangle 4`, `Rounded Rectangle 5`,
    # `TextBox 6` — these are the equation container and equation text
    remove_shape_by_name_substring(slide, [
        "Rounded Rectangle 5",
        "Rectangle 4",
        "TextBox 6",
    ])

    # Insert new equation block as a multi-line equation textbox
    add_equation_textbox(
        slide,
        left=0.8, top=1.4, width=11.7, height=2.0,
        base_size=20,
        lines=[
            [("ℒ", "var"), ("ₜₒₜₐₗ", "var"),  # tota l subscript
             ("(θ) = ", "op"),
             ("ℒ", "var"),
             ("ₚₚₒ", "var"),  # PPO subscript
             (" ( π", "op"), ("θ", "var"),  # π_θ — theta is already standalone
             (" ; v3 dynamics )", "op")],
            [("", "eq"),
             ("            +  λ", "op"), ("ₜₑₘₚ", "var"),
             ("  ·  ‖  ", "op"),
             ("T̂", "var"), ("ₜ₊₁", "var"),
             ("(s, a ; θ", "op"), ("ᵥ₃", "var"),
             (") − ", "op"),
             ("T̂", "var"), ("ₜ₊₁", "var"),
             ("(s, a ; θ", "op"), ("ᵥ₃.₅", "var"),
             (")  ‖²", "op")],
            [("", "eq"),
             ("            +  λ", "op"), ("ₚₒₐᵥₑᵣ", "var"),
             ("  ·  ‖  ", "op"),
             ("P̂", "var"), ("ₜ", "var"),
             ("(s, a ; θ", "op"), ("ᵥ₃", "var"),
             (") − ", "op"),
             ("P̂", "var"), ("ₜ", "var"),
             ("(s, a ; θ", "op"), ("ᵥ₃.₅", "var"),
             (")  ‖²", "op")],
        ],
    )

    # Add a short notation legend under the equation
    add_equation_textbox(
        slide,
        left=0.8, top=3.5, width=11.7, height=0.9,
        base_size=12,
        lines=[
            [("where", "label"),
             ("    θ", "var"), ("ᵥ₃.₅", "var"),
             (" is frozen (calibrated v3.5 weights); ", "label"),
             ("λ", "var"), ("ₜₑₘₚ", "var"),
             (" and ", "label"),
             ("λ", "var"), ("ₚₒₐᵥₑᵣ", "var"),
             (" are the only two hybrid hyperparameters.", "label")],
            [("Forward dynamics use ", "label"),
             ("v3", "var"), (" only. The frozen ", "label"),
             ("v3.5", "var"), (" predictions T̂", "label"),
             ("ₜ₊₁", "var"),
             (" and P̂", "label"), ("ₜ", "var"),
             (" enter only as a soft physical disagreement penalty in the loss.", "label")],
        ],
    )


# ---------------------------------------------------------------------------
# Slide 13 — Surrogate equations (replace table with structured equations)
# ---------------------------------------------------------------------------

def patch_slide_13(prs: Presentation):
    slide = prs.slides[12]
    remove_tables(slide)

    # v3 block
    add_section_header(slide, "v3 — Control-oriented surrogate",
                       left=0.6, top=1.42, width=11.7, color=BLUE)
    add_equation_textbox(
        slide,
        left=0.7, top=1.78, width=11.7, height=0.95,
        base_size=14,
        lines=[
            [("Input :   ", "label"),
             ("x", "var"), ("ₜ", "var"),
             (" = [ ", "op"),
             ("T", "var"), (" ₜₑₒₙₑ", "var"),
             (",  ", "op"),
             ("T", "var"), (" ₐₘₙ", "var"),
             (",  sin ", "op"), ("h", "var"),
             (",  cos ", "op"), ("h", "var"),
             (",  sin ", "op"), ("d", "var"),
             (",  cos ", "op"), ("d", "var"),
             (",  ", "op"),
             ("a", "var"), ("₀", "var"),
             (",  ", "op"),
             ("a", "var"), ("₁", "var"),
             (" ]  ∈ ℝ⁸", "op")],
            [("Output :   ( ", "label"),
             ("T̂", "var"), ("ₜ₊₁", "var"),
             ("ᵥ₃", "var"),
             (",  ", "op"),
             ("P̂", "var"), ("ₜ", "var"),
             ("ᵥ₃", "var"),
             (" )  =  ", "op"),
             ("f", "var"), ("ᵥ₃", "var"),
             (" ( ", "op"),
             ("x", "var"), ("ₜ", "var"),
             (" ; θ", "op"), ("ᵥ₃", "var"),
             (" )", "op")],
        ],
    )

    # v3.5 block
    add_section_header(slide, "v3.5 — Physically informed twin (Stage A/B/C calibrated)",
                       left=0.6, top=2.95, width=11.7, color=SUCCESS)
    add_equation_textbox(
        slide,
        left=0.7, top=3.31, width=11.7, height=1.95,
        base_size=14,
        lines=[
            [("Heat flux :   ", "label"),
             ("q̂", "var"), ("ₜ", "var"),
             ("  =  ", "op"),
             ("g", "var"), ("₆", "var"),
             (" ( ", "op"),
             ("x", "var"), ("ₜ", "var"),
             (" ; θ", "op"), ("₆", "var"),
             (" )", "op")],
            [("Thermal capacitance :   ", "label"),
             ("C", "var"), ("ₜₒₙ", "var"),
             ("  =  ", "op"),
             ("C", "var"), ("ₘᵢₙ", "var"),
             ("  +  10⁵ · softplus( θ", "op"),
             ("C", "var"),
             (" )    such that    ", "op"),
             ("C", "var"), ("ₜₒₙ", "var"),
             ("  ≥  4 × 10⁴  J·K⁻¹", "op")],
            [("Temperature update :   ", "label"),
             ("T̂", "var"), ("ₜ₊₁", "var"),
             ("ᵥ₃.₅", "var"),
             ("  =  clip(  ", "op"),
             ("T", "var"), ("ₜ", "var"),
             ("  +  Δt · ", "op"),
             ("q̂", "var"), ("ₜ", "var"),
             ("  /  ", "op"),
             ("C", "var"), ("ₜₒₙ", "var"),
             (",  15 °C,  35 °C  )", "op")],
            [("Power output :   ", "label"),
             ("P̂", "var"), ("ₜ", "var"),
             ("ᵥ₃.₅", "var"),
             ("  =  ", "op"),
             ("g", "var"), ("ₚ", "var"),
             (" ( ", "op"),
             ("x", "var"), ("ₜ", "var"),
             (" ; θ", "op"), ("ₚ", "var"),
             (" )", "op")],
        ],
    )

    # Hybrid loss block (compact reference)
    add_section_header(slide, "Hybrid training loss (full form on slide 11)",
                       left=0.6, top=5.32, width=11.7, color=ACCENT)
    add_equation_textbox(
        slide,
        left=0.7, top=5.68, width=11.7, height=0.7,
        base_size=14,
        lines=[
            [("ℒ", "var"), ("ₜₒₜₐₗ", "var"),
             ("  =  ", "op"),
             ("ℒ", "var"), ("ₚₚₒ", "var"),
             ("  +  λ", "op"), ("ₜₑₘₚ", "var"),
             (" · ‖ T̂", "op"), ("ₜ₊₁", "var"),
             ("ᵥ₃", "var"),
             (" − T̂", "op"), ("ₜ₊₁", "var"),
             ("ᵥ₃.₅", "var"),
             (" ‖²  +  λ", "op"), ("ₚₒₐᵥₑᵣ", "var"),
             (" · ‖ P̂", "op"), ("ₜ", "var"), ("ᵥ₃", "var"),
             (" − P̂", "op"), ("ₜ", "var"), ("ᵥ₃.₅", "var"),
             (" ‖²", "op")],
        ],
    )


# ---------------------------------------------------------------------------
# Slide 14 — Control interface, reward, and safety metric
# ---------------------------------------------------------------------------

def patch_slide_14(prs: Presentation):
    slide = prs.slides[13]
    remove_tables(slide)

    # Action space
    add_section_header(slide, "Action space",
                       left=0.6, top=1.42, width=11.7, color=BLUE)
    add_equation_textbox(
        slide,
        left=0.7, top=1.78, width=11.7, height=0.95,
        base_size=13,
        lines=[
            [("a", "var"), ("ₜ", "var"),
             ("  =  [ ", "op"),
             ("a", "var"), ("₀", "var"),
             (",  ", "op"),
             ("a", "var"), ("₁", "var"),
             (" ]  ∈  [ −1, +1 ]²", "op")],
            [("T", "var"), ("ₛₜᵤₚₚₗᵧ", "var"),
             ("  =  18 °C  +  ½ · (", "op"),
             ("a", "var"), ("₀", "var"),
             (" + 1) · (35 − 18) °C", "op")],
            [("fan duty cycle  =  clip(  ½ · (", "label"),
             ("a", "var"), ("₁", "var"),
             (" + 1),  0,  1  )", "op")],
        ],
    )

    # Observation
    add_section_header(slide, "17-D observation (MORL canonical)",
                       left=0.6, top=2.85, width=11.7, color=BLUE)
    add_equation_textbox(
        slide,
        left=0.7, top=3.20, width=11.7, height=0.85,
        base_size=13,
        lines=[
            [("o", "var"), ("ₜ", "var"),
             ("  =  [   5 physical features   ⊕   4 cyclic time features   ⊕   5 weather forecasts   ⊕   3 history features   ]", "op")],
            [("                                     (history) :   ", "label"),
             ("a", "var"), ("ₜ₋₁", "var"),
             ("  ⊕  Δ", "op"),
             ("T", "var"), ("ₜ", "var"),
             ("  ⊕  ", "op"),
             ("P", "var"), ("ₜ₋₁", "var")],
        ],
    )

    # Reward
    add_section_header(slide, "Training reward (PPO objective)",
                       left=0.6, top=4.18, width=11.7, color=BLUE)
    add_equation_textbox(
        slide,
        left=0.7, top=4.53, width=11.7, height=0.95,
        base_size=13,
        lines=[
            [("r", "var"), ("ₜ", "var"),
             ("  =  ", "op"),
             ("r", "var"), ("ₜᵣₐₖ", "var"),
             ("  +  ", "op"),
             ("r", "var"), ("ₛₘₒₒₜₕ", "var"),
             ("  +  ", "op"),
             ("r", "var"), ("ₚₒₐᵥₑᵣ", "var"),
             ("  +  ", "op"),
             ("r", "var"), ("ₛₐₒᵥₑ", "var")],
            [("r", "var"), ("ₛₘₒₒₜₕ", "var"),
             ("  =  − 0.05 · ‖ ", "op"),
             ("a", "var"), ("ₜ", "var"),
             ("  −  ", "op"),
             ("a", "var"), ("ₜ₋₁", "var"),
             (" ‖²       ", "op"),
             ("r", "var"), ("ₚₒₐᵥₑᵣ", "var"),
             ("  =  − 3 × 10⁻⁵ · ", "op"),
             ("P", "var"), ("ₜ", "var"),
             ("  (inside comfort band)", "op")],
        ],
    )

    # Safety metric
    add_section_header(slide, "Reported safety metric (BOPTEST-style, independent of training reward)",
                       left=0.6, top=5.62, width=11.7, color=SUCCESS)
    add_equation_textbox(
        slide,
        left=0.7, top=5.97, width=11.7, height=0.85,
        base_size=13,
        lines=[
            [("r", "var"), ("ₜᵢₘₑ", "var"),
             ("  =  mean[  𝟙( ", "op"),
             ("T", "var"), (" ₜₑₒₙₑ", "var"),
             ("  <  21 °C  ∨  ", "op"),
             ("T", "var"), (" ₜₑₒₙₑ", "var"),
             ("  >  24 °C )  ]                  ", "op"),
             ("r", "var"), ("ₛₑᵥ", "var"),
             ("  =  max( max(21 − ", "op"),
             ("T", "var"), ("ₜ", "var"),
             (", 0),  max(", "op"),
             ("T", "var"), ("ₜ", "var"),
             ("  − 24, 0) )", "op")],
            [("m", "var"), ("ₛ", "var"),
             ("  =  ", "op"),
             ("r", "var"), ("ₜᵢₘₑ", "var"),
             ("  +  ", "op"),
             ("r", "var"), ("ₛₑᵥ", "var"),
             ("              ", "op"),
             ("(lower is safer)", "label")],
        ],
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"[INFO] reading  {IN_PATH}")
    prs = Presentation(str(IN_PATH))
    print(f"[INFO] slides   {len(prs.slides)}")
    patch_slide_11(prs)
    print(f"[OK]   patched slide 11 — hybrid loss")
    patch_slide_13(prs)
    print(f"[OK]   patched slide 13 — surrogate equations")
    patch_slide_14(prs)
    print(f"[OK]   patched slide 14 — control interface / reward / m_s")
    prs.save(str(OUT_PATH))
    print(f"[OK]   wrote    {OUT_PATH}")


if __name__ == "__main__":
    main()
