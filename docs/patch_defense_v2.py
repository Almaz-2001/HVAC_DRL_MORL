"""
Second-pass patch for HVAC_DRL_MORL_defense_full.pptx.

Two changes:

1. Simplify formula slides 11, 13, 14.
   The earlier patch used mathematically rigorous typography (hats, theta
   subscripts, parameterised function arguments). For a defense audience
   that is visually dense. This pass replaces those slides with a
   plain-language reading: words first, math second.

2. Insert two new testcase slides covering all four BOPTEST testcases used
   across Blocks 1-3:
     - bestest_air                       (Block 1/2 source)
     - bestest_hydronic_heat_pump        (Block 3 primary)
     - bestest_hydronic                  (Block 3 secondary)
     - singlezone_commercial_hydronic    (Block 3 stretch)

   Slide A: side-by-side comparison table (envelope, heating system,
            actuators, key parameters).
   Slide B: 2x2 visual schematics built from native PowerPoint shapes
            (no external images, no network calls).

   The two new slides are appended at the end of the deck. After running
   this patch, drag them in PowerPoint Slide Sorter view to the position
   recommended by the cover comment (after Speed benchmark, before
   Block 1 divider).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
IN_PATH = ROOT / "docs" / "HVAC_DRL_MORL_defense_full.pptx"
OUT_PATH = IN_PATH

INK = RGBColor(0x1A, 0x1A, 0x2E)
WARN = RGBColor(0xC1, 0x12, 0x1F)
SUCCESS = RGBColor(0x0F, 0x8A, 0x5F)
ACCENT = RGBColor(0xF4, 0xA2, 0x61)
BLUE = RGBColor(0x1E, 0x60, 0x91)
PURPLE = RGBColor(0x5A, 0x4E, 0x7C)
MUTED = RGBColor(0x6B, 0x70, 0x80)
LIGHT_BG = RGBColor(0xF7, 0xF6, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# Generic helpers
# ---------------------------------------------------------------------------

def remove_shapes_after_title(slide, keep_title_count=2, keep_takeaway=False,
                              keep_footer=True):
    """Remove all body shapes; keep title bar + title textbox.

    The current deck builds title bar as the first AUTO_SHAPE and the title
    + subtitle as two TEXT_BOX shapes; we keep those plus optionally the
    takeaway box and the slide footer.
    """
    to_remove = []
    for i, sh in enumerate(slide.shapes):
        # heuristically keep: title bar (rect 0), title textbox (textbox 2),
        # subtitle textbox (textbox 3), divider line (rectangle 4)
        if i <= 3:
            continue
        # keep takeaway box (Rounded Rectangle <last>) if requested
        is_takeaway = False
        if sh.has_text_frame:
            t = sh.text_frame.text or ""
            if t.lstrip().startswith("★"):
                is_takeaway = True
        if is_takeaway and keep_takeaway:
            continue
        # keep footer if requested
        if keep_footer and sh.has_text_frame:
            t = sh.text_frame.text or ""
            if "HVAC DRL/MORL Defense" in t:
                continue
        to_remove.append(sh)
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)


def add_text(slide, text, *, left, top, width, height,
             font_size=14, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_multiline(slide, lines, *, left, top, width, height,
                  font_size=14, color=INK, font="Calibri",
                  align=PP_ALIGN.LEFT, line_spacing=Pt(6)):
    """Add a multi-line textbox; each line is a string."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
    return tb


def add_table(slide, headers, rows, *, left, top, width, height,
              header_fill=INK, header_text=WHITE,
              header_font_size=12, body_font_size=11):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers),
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height)).table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.bold = True; r.font.size = Pt(header_font_size)
        r.font.color.rgb = header_text
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
    for ri, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            cell = tbl.cell(ri, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(body_font_size); r.font.color.rgb = INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
    return tbl


def add_box(slide, *, left, top, width, height, fill=LIGHT_BG, line=INK,
            line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line; shape.line.width = line_width
    return shape


# ---------------------------------------------------------------------------
# 1. Simplify formula slides
# ---------------------------------------------------------------------------

def simplify_slide_11_hybrid_loss(prs: Presentation):
    """Slide 11: simplest plain-language formulation of hybrid loss."""
    s = prs.slides[10]
    remove_shapes_after_title(s, keep_takeaway=True)

    # 1) Plain-language statement (large, central)
    box = add_box(s, left=0.7, top=1.4, width=11.9, height=1.8,
                  fill=LIGHT_BG, line=BLUE, line_width=Pt(1.5))
    add_multiline(s, [
        "Total training loss",
        "",
        "    =   standard PPO loss   (policy rolls out through v3)",
        "    +   λ_temp   ·   (temperature disagreement between v3 and v3.5)²",
        "    +   λ_power  ·   (power disagreement between v3 and v3.5)²",
    ], left=0.9, top=1.5, width=11.5, height=1.6,
       font_size=17, color=INK, font="Calibri")

    # 2) Variable table
    add_table(s,
        headers=["Symbol", "Meaning", "Value"],
        rows=[
            ("λ_temp",  "Weight on temperature-disagreement penalty",
             "0.10  (thermostatic) ; 0.00  (HDRL, MORL)"),
            ("λ_power", "Weight on power-disagreement penalty",
             "5 · 10⁻⁵  (universal across controllers)"),
            ("v3",      "Control-oriented neural surrogate (used for rollouts)",
             "trained on direct-T_supply trajectories"),
            ("v3.5",    "Physically informed twin (FROZEN; used only in loss)",
             "Stage A/B/C calibrated, with explicit C_zon"),
        ],
        left=0.7, top=3.4, width=11.9, height=2.4,
        header_font_size=12, body_font_size=12)

    # The existing takeaway box is preserved by remove_shapes_after_title


def simplify_slide_13_surrogate_equations(prs: Presentation):
    """Slide 13: replace dense equations with two concise concept blocks."""
    s = prs.slides[12]
    remove_shapes_after_title(s, keep_takeaway=True)

    # v3 block
    add_box(s, left=0.6, top=1.4, width=11.9, height=1.8,
            fill=LIGHT_BG, line=BLUE, line_width=Pt(1.5))
    add_text(s, "v3  —  control-oriented surrogate",
             left=0.85, top=1.5, width=11.5, height=0.4,
             font_size=15, bold=True, color=BLUE)
    add_multiline(s, [
        "•   Input  :   8 features  =  zone & ambient temperature, "
        "time-of-day (sin/cos), day-of-year (sin/cos), and 2 action commands",
        "•   Output :   predicted next-step zone temperature  +  predicted "
        "total HVAC power",
        "•   Type   :   pure feed-forward neural network  (no explicit physics)",
        "•   Role   :   smooth, fast rollout environment for PPO",
    ], left=0.85, top=1.9, width=11.5, height=1.2,
       font_size=13, color=INK)

    # v3.5 block
    add_box(s, left=0.6, top=3.4, width=11.9, height=2.4,
            fill=LIGHT_BG, line=SUCCESS, line_width=Pt(1.5))
    add_text(s, "v3.5  —  physically informed twin",
             left=0.85, top=3.5, width=11.5, height=0.4,
             font_size=15, bold=True, color=SUCCESS)
    add_multiline(s, [
        "•   Same 8 inputs as v3",
        "•   Adds an explicit physical parameter:  C_zon = zone thermal "
        "capacitance, in J/K  (identifiable from data)",
        "•   Temperature update is hand-written physics:",
        "         T_next   =   T_now   +   Δt  ·  (heat flux from NN)  /  C_zon",
        "•   Power is still predicted by a neural head, but constrained by "
        "Stage C calibration",
        "•   Type   :   grey-box   (NN for heat flux + physics for temperature)",
    ], left=0.85, top=3.9, width=11.5, height=1.85,
       font_size=13, color=INK)


def simplify_slide_14_control(prs: Presentation):
    """Slide 14: simplified control interface / reward / safety."""
    s = prs.slides[13]
    remove_shapes_after_title(s, keep_takeaway=True)

    # action
    add_box(s, left=0.6, top=1.4, width=5.85, height=1.5,
            fill=LIGHT_BG, line=BLUE)
    add_text(s, "What the controller emits",
             left=0.75, top=1.5, width=5.6, height=0.35,
             font_size=14, bold=True, color=BLUE)
    add_multiline(s, [
        "•   2 numbers in [-1, +1]  per step",
        "•   Mapped to supply temperature (18-35 °C)",
        "•   and to fan speed (0-1)",
    ], left=0.75, top=1.85, width=5.6, height=1.0,
       font_size=12, color=INK)

    # observation
    add_box(s, left=6.7, top=1.4, width=5.85, height=1.5,
            fill=LIGHT_BG, line=BLUE)
    add_text(s, "What the controller sees",
             left=6.85, top=1.5, width=5.6, height=0.35,
             font_size=14, bold=True, color=BLUE)
    add_multiline(s, [
        "•   17 features  :   5 physical  +  4 time  +",
        "    5 weather forecasts  +  3 history",
        "•   Rich observation reduces need for λ_temp anchor",
    ], left=6.85, top=1.85, width=5.6, height=1.0,
       font_size=12, color=INK)

    # reward
    add_box(s, left=0.6, top=3.05, width=11.95, height=1.5,
            fill=LIGHT_BG, line=ACCENT)
    add_text(s, "What the controller optimises during training (PPO reward)",
             left=0.75, top=3.15, width=11.7, height=0.35,
             font_size=14, bold=True, color=ACCENT)
    add_multiline(s, [
        "Reward  =  comfort tracking  +  smoothness  +  energy  +  hybrid disagreement penalty",
        "•   comfort tracking   :  inside band [21 °C, 24 °C]  →  positive ; outside  →  negative penalty",
        "•   smoothness         :  − 0.05  ·  (change in action between consecutive steps)²",
        "•   energy             :  − 3 · 10⁻⁵  ·  total power, but only when temperature is already in band",
    ], left=0.75, top=3.50, width=11.7, height=1.1,
       font_size=12, color=INK)

    # safety metric
    add_box(s, left=0.6, top=4.7, width=11.95, height=1.4,
            fill=LIGHT_BG, line=SUCCESS)
    add_text(s, "What the paper reports as the verdict  :  m_s safety metric (BOPTEST-style)",
             left=0.75, top=4.8, width=11.7, height=0.35,
             font_size=14, bold=True, color=SUCCESS)
    add_multiline(s, [
        "•   r_time  =  fraction of time the zone temperature is outside the comfort band",
        "•   r_sev   =  maximum °C excursion outside the band  (worst single violation)",
        "•   m_s     =  r_time  +  r_sev          lower is safer ;  PI baseline m_s used as the normaliser",
    ], left=0.75, top=5.15, width=11.7, height=0.95,
       font_size=12, color=INK)


# ---------------------------------------------------------------------------
# 2. Testcase slides
# ---------------------------------------------------------------------------

def add_title_bar(s, title, subtitle):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                             Inches(13.333), Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    add_text(s, title, left=0.4, top=0.18, width=12.5, height=0.65,
             font_size=24, bold=True, color=INK)
    add_text(s, subtitle, left=0.4, top=0.78, width=12.5, height=0.5,
             font_size=13, italic=True, color=MUTED)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.20),
                              Inches(12.5), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = MUTED
    line.line.fill.background()


def add_takeaway(s, text, left=0.4, top=6.5, width=12.5, height=0.55):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(left), Inches(top),
                             Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = SUCCESS; box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10); tf.margin_right = Pt(10)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "★ Takeaway:  "
    r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = SUCCESS
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(11); r2.font.color.rgb = INK


def add_footer(s, label):
    add_text(s, label, left=0.4, top=7.1, width=12.5, height=0.3,
             font_size=9, italic=True, color=MUTED)


def insert_testcase_overview(prs: Presentation):
    """New slide: side-by-side testcase comparison table."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s,
        "BOPTEST testcases used in this study",
        "All four testcases come from the IBPSA Project 1 BOPTEST suite "
        "(github.com/ibpsa/project1-boptest). Sources of building physics, "
        "HVAC schematics, and KPI definitions: BOPTEST 1.0 documentation.")
    rows = [
        ("Source testcase", "bestest_air",
         "ASHRAE BESTEST single-zone envelope; "
         "forced-air heating",
         "Direct supply-temperature override + fan command",
         "C_zon ≈ 4.41 · 10⁵ J/K (identified)",
         "Blocks 1 and 2"),
        ("Block 3 primary", "bestest_hydronic_heat_pump",
         "Same BESTEST envelope; hydronic loop with air-source heat pump",
         "Setpoint + on/off heat pump + pump + fan overrides",
         "C_zon ≈ 8.35 · 10⁵ J/K (1.89× source)",
         "Block 3, mode=none/partial/full"),
        ("Block 3 secondary", "bestest_hydronic",
         "Same BESTEST envelope; hydronic loop with boiler and radiators",
         "Setpoint + boiler modulation + circulation pump",
         "C_zon ≈ 8.62 · 10⁵ J/K (1.95× source)",
         "Block 3, mode=none/full"),
        ("Block 3 stretch", "singlezone_commercial_hydronic",
         "Larger commercial-scale single-zone envelope; hydronic distribution",
         "Setpoint + circulation pump + boiler/heat-source modulation",
         "C_zon ≈ 8.43 · 10⁵ J/K (1.91× source)",
         "Block 3, mode=none/full"),
    ]
    add_table(s,
        headers=["Role in study", "BOPTEST id", "Envelope and HVAC",
                 "Actuator interface", "Thermal capacitance C_zon",
                 "Where in this study"],
        rows=rows,
        left=0.3, top=1.45, width=12.7, height=4.85,
        header_font_size=10, body_font_size=9.5)

    add_takeaway(s,
        "Four testcases, one envelope family at residential scale plus one larger commercial-scale variant. "
        "Hydronic family shows a stable C_zon ratio ≈ 1.9× the bestest_air source value across all three Block 3 testcases.")
    add_footer(s, "HVAC DRL/MORL Defense  ·  Testcase overview")


def draw_envelope(s, *, left, top, width, height, title, subtitle, color,
                  heat_source_label, distribution_label, actuator_lines):
    """Draw a small testcase schematic into a sub-region of the slide."""
    # outer card
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(left), Inches(top),
                              Inches(width), Inches(height))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = color; card.line.width = Pt(1.5)

    # title strip
    add_text(s, title, left=left + 0.15, top=top + 0.12,
             width=width - 0.3, height=0.32,
             font_size=13, bold=True, color=color)
    add_text(s, subtitle, left=left + 0.15, top=top + 0.45,
             width=width - 0.3, height=0.28,
             font_size=9, italic=True, color=MUTED)

    # building envelope rectangle (centered)
    env_w = width - 2.0
    env_h = 1.10
    env_left = left + (width - env_w) / 2
    env_top = top + 0.85
    env = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(env_left), Inches(env_top),
                             Inches(env_w), Inches(env_h))
    env.fill.solid(); env.fill.fore_color.rgb = WHITE
    env.line.color.rgb = INK; env.line.width = Pt(1.2)
    add_text(s, "Zone  (T_zone)", left=env_left + 0.1, top=env_top + 0.05,
             width=env_w - 0.2, height=0.30,
             font_size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, distribution_label,
             left=env_left + 0.1, top=env_top + 0.40,
             width=env_w - 0.2, height=0.55,
             font_size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # heat source box (left of envelope)
    hs_w = 0.85
    hs_h = 0.55
    hs_left = left + 0.18
    hs_top = env_top + (env_h - hs_h) / 2
    hs = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(hs_left), Inches(hs_top),
                            Inches(hs_w), Inches(hs_h))
    hs.fill.solid(); hs.fill.fore_color.rgb = color
    hs.line.color.rgb = INK; hs.line.width = Pt(0.8)
    add_text(s, heat_source_label, left=hs_left, top=hs_top,
             width=hs_w, height=hs_h,
             font_size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # arrow from heat source to envelope
    arrow = s.shapes.add_connector(2,  # straight connector
                                    Inches(hs_left + hs_w),
                                    Inches(hs_top + hs_h / 2),
                                    Inches(env_left),
                                    Inches(env_top + env_h / 2))
    arrow.line.color.rgb = color
    arrow.line.width = Pt(2.0)
    # arrowhead via XML
    ln = arrow.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"),
                          {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)

    # actuator interface list under envelope
    add_text(s, "Overrides:",
             left=env_left, top=env_top + env_h + 0.10,
             width=env_w, height=0.25,
             font_size=10, bold=True, color=INK)
    add_multiline(s, actuator_lines,
                  left=env_left, top=env_top + env_h + 0.40,
                  width=env_w, height=0.85,
                  font_size=9, color=INK,
                  line_spacing=Pt(2))


def insert_testcase_schematics(prs: Presentation):
    """New slide: 2x2 schematic grid of the four testcases."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s,
        "Testcase schematics",
        "How heat reaches the zone in each of the four BOPTEST testcases. "
        "Simplified diagrams drawn from BOPTEST 1.0 testcase documentation; "
        "see github.com/ibpsa/project1-boptest for full Modelica models.")

    # 2x2 grid
    cells = [
        # (row, col, ...)
        (0, 0,
            "bestest_air",
            "Source testcase  (Blocks 1 and 2)",
            BLUE,
            "Air-side\nheater",
            "Forced-air supply duct  →  zone",
            ["•   oveTSetHea_u   (setpoint)",
             "•   con_oveTSetHea_activate",
             "•   fcu_oveFan_u   (fan duty)"]),
        (0, 1,
            "bestest_hydronic_heat_pump",
            "Block 3 primary  (closest neighbour)",
            SUCCESS,
            "Air-source\nheat pump",
            "Hydronic loop  →  fan coil  →  zone",
            ["•   oveTSet_u   (zone setpoint)",
             "•   oveHeaPumY_u  (HP on/off)",
             "•   ovePum_u, oveFan_u"]),
        (1, 0,
            "bestest_hydronic",
            "Block 3 secondary  (mid difficulty)",
            ACCENT,
            "Boiler\n(gas/elec)",
            "Hydronic loop  →  radiators  →  zone",
            ["•   oveTSet_u  (setpoint)",
             "•   ove modulation (boiler)",
             "•   ovePum_u  (circulation)"]),
        (1, 1,
            "singlezone_commercial_hydronic",
            "Block 3 stretch  (larger envelope)",
            PURPLE,
            "Commercial\nheat source",
            "Hydronic distribution  →  larger zone",
            ["•   oveTSet_u  (setpoint)",
             "•   ovePum_u, ove modulation",
             "•   Building scale × ~larger"]),
    ]
    col_w = 6.2
    row_h = 2.55
    base_left = 0.45
    base_top = 1.4
    for row, col, title, sub, color, hs_label, dist_label, acts in cells:
        draw_envelope(s,
                      left=base_left + col * (col_w + 0.25),
                      top=base_top + row * (row_h + 0.05),
                      width=col_w, height=row_h,
                      title=title, subtitle=sub, color=color,
                      heat_source_label=hs_label,
                      distribution_label=dist_label,
                      actuator_lines=acts)

    add_takeaway(s,
        "All four testcases share the same single-zone envelope family at residential scale, "
        "but differ in heat source type and actuator interface. This is why Block 3 cannot be "
        "a literal direct-T_supply transfer; it is adapter-mediated.")
    add_footer(s, "HVAC DRL/MORL Defense  ·  Testcase schematics")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"[INFO] reading  {IN_PATH}")
    prs = Presentation(str(IN_PATH))
    n_before = len(prs.slides)
    print(f"[INFO] slides before: {n_before}")

    simplify_slide_11_hybrid_loss(prs)
    print(f"[OK]   simplified slide 11 (hybrid loss)")
    simplify_slide_13_surrogate_equations(prs)
    print(f"[OK]   simplified slide 13 (surrogate equations)")
    simplify_slide_14_control(prs)
    print(f"[OK]   simplified slide 14 (control / reward / m_s)")

    insert_testcase_overview(prs)
    print(f"[OK]   appended testcase overview slide  (NEW)")
    insert_testcase_schematics(prs)
    print(f"[OK]   appended testcase schematics slide (NEW)")

    n_after = len(prs.slides)
    prs.save(str(OUT_PATH))
    print(f"[INFO] slides after:  {n_after}  (+{n_after - n_before} new)")
    print(f"[OK]   wrote          {OUT_PATH}")
    print()
    print("Manual step in PowerPoint:")
    print("  1. Open the file in PowerPoint.")
    print("  2. View -> Slide Sorter.")
    print("  3. Drag the two new appended slides to position AFTER slide 16 "
          "(Speed benchmark) and BEFORE slide 17 (Block 1 divider).")


if __name__ == "__main__":
    main()
