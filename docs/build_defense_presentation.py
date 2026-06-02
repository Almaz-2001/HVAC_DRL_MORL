"""
Build the full HVAC DRL/MORL defense presentation.

Reads all article-real figures from reports/figures/article_real/ and
selected auxiliary figures from reports/figures/, and assembles a 16:9
PPTX deck covering Blocks 1, 2, 3 with the same numerical evidence used
in the manuscript.

Two slides remain intentionally empty for the author to populate later:
- Related Works
- Literature Review

Output:
    docs/HVAC_DRL_MORL_defense_concise.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
FIG_REAL = ROOT / "reports" / "figures" / "article_real"
FIG_AUX = ROOT / "reports" / "figures"
OUT_PATH = ROOT / "docs" / "HVAC_DRL_MORL_defense_concise.pptx"

# ----- palette (consistent with paper figures) -----
INK = RGBColor(0x1A, 0x1A, 0x2E)
WARN = RGBColor(0xC1, 0x12, 0x1F)
SUCCESS = RGBColor(0x0F, 0x8A, 0x5F)
ACCENT = RGBColor(0xF4, 0xA2, 0x61)
BLUE = RGBColor(0x1E, 0x60, 0x91)
MUTED = RGBColor(0x6B, 0x70, 0x80)
LIGHT_BG = RGBColor(0xF7, 0xF6, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def make_presentation() -> Presentation:
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)


def add_text(slide, text, *, left, top, width, height,
             font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_title_bar(slide, title, subtitle=None):
    # narrow accent stripe at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  Inches(13.333), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    # title
    add_text(slide, title, left=0.4, top=0.18, width=12.5, height=0.65,
             font_size=26, bold=True, color=INK)
    if subtitle:
        add_text(slide, subtitle, left=0.4, top=0.78, width=12.5, height=0.5,
                 font_size=14, italic=True, color=MUTED)
    # footer line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.4), Inches(1.20),
                                   Inches(12.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = MUTED
    line.line.fill.background()


def add_footer(slide, idx, total, section=""):
    footer_text = f"HVAC DRL/MORL Defense | Slide {idx}/{total}"
    if section:
        footer_text += f" | {section}"
    add_text(slide, footer_text, left=0.4, top=7.1, width=12.5, height=0.3,
             font_size=9, color=MUTED, italic=True)


def add_image_centered(slide, image_path: Path, *, top, max_height=4.9, max_width=12.5):
    if not image_path.exists():
        # placeholder rectangle if image missing
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(2), Inches(top),
                                     Inches(max_width - 4), Inches(max_height))
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_BG
        ph.line.color.rgb = WARN
        ph.line.width = Pt(1.5)
        tf = ph.text_frame
        tf.text = f"[MISSING: {image_path.name}]"
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WARN
                r.font.italic = True
        return
    # try to size to fit; let pptx infer aspect by setting only height
    pic = slide.shapes.add_picture(str(image_path), Inches(0.4), Inches(top),
                                    height=Inches(max_height))
    # center horizontally if too wide
    slide_w = Inches(13.333)
    if pic.width > Inches(max_width):
        # rescale
        ratio = pic.width / pic.height
        pic.height = Inches(max_height)
        pic.width = int(pic.height * ratio)
    pic.left = int((slide_w - pic.width) / 2)


def add_takeaway_box(slide, text, *, left=0.4, top=5.95, width=12.5, height=0.9):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = SUCCESS
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Takeaway: "
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = SUCCESS
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(13)
    r2.font.color.rgb = INK


def add_graph_note(slide, text, *, top=1.27):
    """Short factual note explaining how to read a chart."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.65), Inches(top),
                                  Inches(12.05), Inches(0.46))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEF, 0xF5, 0xF2)
    box.line.color.rgb = SUCCESS
    box.line.width = Pt(0.8)
    add_text(slide, text, left=0.85, top=top + 0.06, width=11.65, height=0.34,
             font_size=10.5, color=INK)
    return box


def add_bullets(slide, items, *, left, top, width, height,
                font_size=16, color=INK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.name = "Calibri"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color


def add_table(slide, headers, rows, *, left, top, width, height,
              header_fill=INK, header_text=WHITE):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers),
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height)).table
    # headers
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = Pt(12)
        r.font.color.rgb = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
    # rows
    for ri, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            cell = tbl.cell(ri, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(v)
            r.font.size = Pt(11)
            r.font.color.rgb = INK
            # zebra
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE


# ===========================================================================
# Slide constructors
# ===========================================================================

SLIDES_TOTAL = 40  # updated dynamically below


def slide_title(prs, idx):
    s = blank_slide(prs)
    # background block
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                            prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()
    # title text
    add_text(s, "When Predictive Surrogates Fail as RL Environments",
             left=0.8, top=1.8, width=12, height=1.2,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, "A Calibrated Physical Twin as Soft Regularizer for HVAC Control",
             left=0.8, top=3.05, width=12, height=0.7,
             font_size=22, italic=True, color=ACCENT, align=PP_ALIGN.LEFT)
    # accent rule
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.8), Inches(3.9),
                               Inches(2.5), Inches(0.05))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    # author block
    add_text(s, "Almaz Sapargali", left=0.8, top=4.3, width=12, height=0.5,
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, "HVAC DRL / MORL Project Defense", left=0.8, top=4.85,
             width=12, height=0.4, font_size=14, color=MUTED, align=PP_ALIGN.LEFT)
    add_text(s, "BOPTEST • bestest_air, bestest_hydronic family",
             left=0.8, top=5.25, width=12, height=0.4,
             font_size=14, color=MUTED, align=PP_ALIGN.LEFT)
    # date / version
    add_text(s, "Target journal: Results in Engineering (Elsevier, Q1)",
             left=0.8, top=6.3, width=12, height=0.4,
             font_size=12, italic=True, color=ACCENT, align=PP_ALIGN.LEFT)


def slide_outline(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Outline",
                  subtitle="Three blocks, one pre-registered audit chain")
    items_left = [
        "1. Relevance, novelty, research questions",
        "2. Methodology overview",
        "3. Speed benchmark (85x over BOPTEST RTE)",
        "4. Block 1 — Surrogate Fidelity",
        "5. Block 2 — Hybrid Construction & Control",
    ]
    items_right = [
        "6. Block 3 — Transferability",
        "7. Pre-registration audit trail",
        "8. Hou-and-Evins compliance",
        "9. Limitations and threats to validity",
        "10. Conclusion and Block 4 future work",
    ]
    add_bullets(s, items_left, left=0.6, top=1.6, width=6.0, height=5.0, font_size=18)
    add_bullets(s, items_right, left=6.8, top=1.6, width=6.0, height=5.0, font_size=18)
    add_footer(s, idx, SLIDES_TOTAL, "Outline")


def slide_motivation(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Motivation",
                  subtitle="Why reinforcement learning for HVAC needs a different surrogate")
    items = [
        "Buildings account for ≈ 40% of global energy use; HVAC ≈ 50% of that.",
        "Reinforcement learning is sample-expensive: training directly in BOPTEST RTE HTTP costs ~21 env-steps/s.",
        "Standard practice: train on a neural surrogate, deploy on live BOPTEST.",
        "Implicit assumption: a more physically faithful surrogate ⇒ a better RL environment.",
        "This work falsifies that assumption and proposes a hybrid resolution.",
    ]
    add_bullets(s, items, left=0.6, top=1.6, width=12.0, height=4.5, font_size=18)
    add_takeaway_box(s, "Calibration and RL utility are not the same goal — they trade off.")
    add_footer(s, idx, SLIDES_TOTAL, "Motivation")


def slide_relevance(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Relevance / actuality",
                  subtitle="Why the problem matters now")
    rows = [
        ("Energy and carbon pressure",
         "Buildings are a major energy consumer; HVAC is the controllable load where operational optimization can matter immediately."),
        ("RL training bottleneck",
         "Live BOPTEST-style simulators are too slow for PPO-scale training; the project measured only ~21 env-steps/s through the HTTP-Docker loop."),
        ("Surrogate risk",
         "Most HVAC-RL workflows rely on surrogates, but predictive fidelity alone does not guarantee closed-loop control utility."),
        ("Deployment risk",
         "A controller that looks good on a native testcase can fail after actuator-interface changes; Block 3 tests that transfer boundary explicitly."),
    ]
    add_table(s, headers=["Relevance axis", "Concrete reason in this project"],
              rows=rows, left=0.7, top=1.5, width=12.0, height=3.2)
    add_bullets(s, [
        "The practical question is not only 'can we predict temperature?' but 'can a surrogate train a controller that survives live BOPTEST?'",
        "That distinction is the core reason for separating Block 1 predictive validity, Block 2 control utility, and Block 3 transferability.",
    ], left=0.8, top=5.0, width=11.7, height=1.1, font_size=14)
    add_takeaway_box(s,
        "The project is relevant because it tests the exact failure mode that blocks surrogate-trained HVAC RL from becoming deployable.")
    add_footer(s, idx, SLIDES_TOTAL, "Motivation")


def slide_research_questions(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Research questions",
                  subtitle="One falsifiable question per block")
    rows = [
        ("Block 1", "Surrogate Fidelity",
         "Does explicit C_zon identification (Stage A/B/C) improve predictive validity over a black-box surrogate?"),
        ("Block 2", "Control Utility",
         "Does the calibrated physical twin make a better RL training environment than the smooth v3 surrogate?"),
        ("Block 3", "Transferability",
         "Does the hybrid recipe transfer to related BOPTEST testcases under pre-registered recalibration regimes?"),
    ]
    add_table(s,
              headers=["Block", "Theme", "Falsifiable question"],
              rows=rows,
              left=0.6, top=1.6, width=12.0, height=2.5)
    add_takeaway_box(s,
        "Three separate hypotheses, three pre-registered protocols, four audit-anchor commits in git.")
    add_footer(s, idx, SLIDES_TOTAL, "Research questions")


def slide_scientific_novelty(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Scientific novelty",
                  subtitle="What is new relative to a standard surrogate + PPO study")
    rows = [
        ("1", "Fidelity-to-RL gap isolated",
         "Shows that calibrated predictive validity can fail as a standalone RL environment: v3.5 has strong rollout RMSE but >4 C live transfer RMSE."),
        ("2", "Hybrid role split",
         "Uses v3.5 as a frozen soft regularizer while rollouts stay on smooth v3 dynamics; physics is in the loss, not the environment."),
        ("3", "Controller-family specificity",
         "Demonstrates that the optimal disagreement weight is architecture-dependent: thermostatic accepts lambda_temp=0.10; HDRL and MORL reject it."),
        ("4", "Pre-registered falsification trail",
         "Reports failed hypotheses explicitly: MORL action-saturation hypothesis falsified at N=5; Block 3 controller transfer is not deployment-ready."),
        ("5", "Component-level transferability",
         "Separates surrogate transfer from controller transfer across three hydronic testcases; surrogate transfers, frozen controller does not."),
    ]
    add_table(s, headers=["#", "Novelty", "Evidence"],
              rows=rows, left=0.5, top=1.35, width=12.35, height=4.8)
    add_takeaway_box(s,
        "The contribution is not a single better controller; it is a falsifiable decomposition of predictive fidelity, control utility, and transferability.")
    add_footer(s, idx, SLIDES_TOTAL, "Scientific novelty")


def slide_empty_placeholder(prs, idx, title, hint):
    s = blank_slide(prs)
    add_title_bar(s, title, subtitle="[to be added by author]")
    # placeholder content area
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(1.5), Inches(2.0),
                              Inches(10.3), Inches(4.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = MUTED
    box.line.width = Pt(1)
    box.line.dash_style = 2  # dashed
    add_text(s, hint, left=2.0, top=2.5, width=9.3, height=3.0,
             font_size=18, italic=True, color=MUTED, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, idx, SLIDES_TOTAL, title)


def slide_methodology_overview(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Methodology — three components",
                  subtitle="v3 + v3.5 + hybrid loss; one shared evaluation protocol")
    rows = [
        ("v3 surrogate",
         "Comfort-oriented, direct-T_supply, two-headed neural model. Smooth dynamics suitable for PPO rollouts."),
        ("v3.5 surrogate",
         "Physically informed grey-box with explicit C_zon. Inverse-calibrated through Stage A (telemetry preprocessing), Stage B (C_zon identification), Stage C (residual heads)."),
        ("Hybrid backend",
         "Policy evolves under v3 dynamics; calibrated v3.5 enters the loss as a frozen physical disagreement penalty (λ_temp · ‖T_v3 − T_v3.5‖² + λ_power · ‖P_v3 − P_v3.5‖²)."),
        ("Controller families",
         "PI baseline (BOPTEST built-in), thermostatic PPO, HDRL (hierarchical), 17-D MORL (preference-conditioned)."),
        ("Evaluation",
         "BOPTEST RTE live; Δt = 900 s; comfort band [21, 24] °C; per-window for in-testcase + yearly for cross-testcase transfer."),
    ]
    add_table(s,
              headers=["Component", "Description"],
              rows=rows,
              left=0.6, top=1.45, width=12.0, height=4.3)
    add_takeaway_box(s,
        "Calibrated v3.5 is in the loss, never in the forward pass. This is the structural decoupling that resolves the fidelity-to-RL gap.")
    add_footer(s, idx, SLIDES_TOTAL, "Methodology")


def slide_pipeline_schematic(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Pipeline schematic",
                  subtitle="Data → calibration → hybrid backend → controllers → live BOPTEST")
    add_image_centered(s, FIG_REAL / "main_fig1_pipeline_schematic.png",
                       top=1.45, max_height=5.4)
    add_footer(s, idx, SLIDES_TOTAL, "Methodology")


def slide_hybrid_loss(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Hybrid loss",
                  subtitle="The single equation that defines the contribution")
    # the equation as a centered text block
    eq_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1.5), Inches(2.0),
                                 Inches(10.3), Inches(1.8))
    eq_box.fill.solid()
    eq_box.fill.fore_color.rgb = LIGHT_BG
    eq_box.line.color.rgb = INK
    eq_box.line.width = Pt(1.2)
    add_text(s,
        "L_total  =  L_PPO( π_θ ; v3 )  +  λ_temp · ‖ T_v3(s,a) − T_v3.5(s,a) ‖²  +  λ_power · ‖ P_v3(s,a) − P_v3.5(s,a) ‖²",
        left=1.7, top=2.4, width=9.9, height=1.0,
        font_size=20, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    items = [
        "Policy rolls out through v3 dynamics. Smooth, fast, control-friendly.",
        "Calibrated v3.5 is frozen. It contributes only the disagreement penalty.",
        "Only two extra hyperparameters: λ_temp and λ_power.",
        "Decouples training environment (v3) from physical anchor (v3.5).",
    ]
    add_bullets(s, items, left=1.5, top=4.2, width=10.3, height=2.2, font_size=15)
    add_takeaway_box(s,
        "λ values turn out to depend on the controller family — this becomes a finding, not an arbitrary tuning choice.")
    add_footer(s, idx, SLIDES_TOTAL, "Methodology")


def slide_surrogate_architecture(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Surrogate architectures",
                  subtitle="What exactly is different between v3, v3.5, and the hybrid backend")
    rows = [
        ("Input vector x_t",
         "8 features: T_zone, T_amb, sin/cos hour, sin/cos day, a0, a1"),
        ("v3 heat branch",
         "HeatFlowNetV2: Linear(8,64)-LayerNorm-Tanh, residual 64-block, Linear(64,32)-Tanh, Linear(32,1); 7,105 params"),
        ("v3 power branch",
         "PowerNetV2: Linear(8,32)-Tanh-Linear(32,32)-Tanh-Linear(32,1)-Softplus; 1,377 params"),
        ("v3.5 change",
         "Same branch topology, but heat branch predicts q_net and temperature evolves through explicit C_zon; total 8,483 params"),
        ("Hybrid backend",
         "Forward dynamics use v3; frozen calibrated v3.5 is evaluated in parallel only to compute disagreement penalties"),
    ]
    add_table(s, headers=["Element", "Implementation detail"],
              rows=rows, left=0.45, top=1.35, width=12.45, height=4.6)
    add_takeaway_box(s,
        "The hybrid is not a larger black-box model. It is a two-model training objective: smooth v3 dynamics plus a frozen physical censor.")
    add_footer(s, idx, SLIDES_TOTAL, "Methodology")


def slide_surrogate_equations(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Surrogate equations",
                  subtitle="Readable mathematical form of the two surrogate families")
    rows = [
        ("v3 control surrogate",
         "x_t = [T_z, T_amb, sin h, cos h, sin d, cos d, a0, a1]\n"
         "(T_{t+1}^{v3}, P_t^{v3}) = f_{v3}(x_t)"),
        ("v3.5 physical twin",
         "q_t = g_q(x_t)\n"
         "C_zon = C_min + 1e5 * softplus(theta_C)\n"
         "T_{t+1}^{v35} = clip(T_t + Delta_t * q_t / C_zon, 15, 35)\n"
         "P_t^{v35} = g_P(x_t)"),
        ("Hybrid training loss",
         "L_total = L_PPO(pi_theta; v3)\n"
         "        + lambda_T * |T_{t+1}^{v3} - T_{t+1}^{v35}|^2\n"
         "        + lambda_P * |P_t^{v3} - P_t^{v35}|^2"),
    ]
    add_table(s, headers=["Concept", "Formula"],
              rows=rows, left=0.65, top=1.35, width=12.05, height=4.6)
    add_takeaway_box(s,
        "v3.5 changes the temperature update equation, not just the network size: temperature is mediated by an identifiable thermal capacitance.")
    add_footer(s, idx, SLIDES_TOTAL, "Methodology")


def slide_control_interface_reward(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Control interface, reward, and safety metric",
                  subtitle="The exact signals optimized by PPO and reported in BOPTEST")
    rows = [
        ("Action space",
         "a = [a0, a1] in [-1,1]^2\n"
         "T_supply = 18 + 0.5*(a0+1)*(35-18) C\n"
         "fan = clip(0.5*(a1+1), 0, 1)"),
        ("17-D observation",
         "5 physical features + 4 cyclic time features + 5 weather forecasts + 3 history features\n"
         "history = previous action(2) + delta-T feature(1)"),
        ("Thermostatic reward",
         "r = r_track + r_smooth + r_power + r_disagree\n"
         "r_smooth = -0.05 * ||a_t - a_{t-1}||^2\n"
         "r_power = -3e-5 * P_t only when inside comfort band\n"
         "r_disagree = -lambda_T * dT - lambda_P * dP"),
        ("Reported safety metric",
         "r_time = mean(1[T_z < 21 or T_z > 24])\n"
         "r_sev = max(max(21-T_z), max(T_z-24))\n"
         "m_s = r_time + r_sev; lower is safer"),
    ]
    add_table(s, headers=["Object", "Definition"],
              rows=rows, left=0.45, top=1.35, width=12.45, height=4.75)
    add_takeaway_box(s,
        "The training reward is shaped for learning; the paper verdict uses the independent BOPTEST-style m_s safety metric.")
    add_footer(s, idx, SLIDES_TOTAL, "Methodology")


def slide_controller_nn_details(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Controller neural-network details",
                  subtitle="PPO policies, controller families, and why architectures behave differently")
    rows = [
        ("Thermostatic PPO",
         "Single continuous PPO policy; MlpPolicy with [256,256] or weather-GRU extractor when forecast features are active; n_steps=1024, gamma=0.99."),
        ("HDRL",
         "Two PPO policies: winter agent + summer agent; each MlpPolicy [256,256], n_steps=1024, batch_size=2048; high/low seasonal split."),
        ("MORL",
         "Preference-conditioned PPO on 17-D TSup-style observation; MlpPolicy, learning_rate=3e-4, n_steps=2048, batch_size=64, gamma=0.99."),
        ("MORL weights",
         "w = (w_comfort, w_energy, w_safety), normalized to sum to 1; Pareto sweep uses five comfort/energy points; canonical N=5 uses neutral and practical points."),
        ("Why lambda differs",
         "Thermostatic PPO benefits from lambda_T=0.10; HDRL and MORL reject lambda_T>0 because hierarchy/rich observation already constrain temperature behavior."),
    ]
    add_table(s, headers=["Controller", "Network / training detail"],
              rows=rows, left=0.45, top=1.35, width=12.45, height=4.65)
    add_takeaway_box(s,
        "The controller-family result is architectural: the same physical regularizer interacts differently with flat, hierarchical, and preference-conditioned policies.")
    add_footer(s, idx, SLIDES_TOTAL, "Methodology")


def slide_speed_benchmark(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Speed benchmark",
                  subtitle="Why an in-process surrogate is needed for RL training")
    rows = [
        ("BOPTEST RTE HTTP–Docker", "21.0", "1.0x", "Standard production deployment"),
        ("v3 surrogate", "4 626", "220.2x", "Smoothest, fastest"),
        ("v3.5 calibrated surrogate", "2 400", "114.2x", "Physical twin alone"),
        ("hybrid_l010 (canonical)", "1 787", "85.0x", "Used for all Block 2/3 training"),
    ]
    add_table(s,
              headers=["Backend", "env-steps/s", "Speed-up", "Notes"],
              rows=rows,
              left=0.8, top=1.7, width=11.7, height=2.4)
    items = [
        "Same control protocol (Δt = 900 s, 100 episodes × 96 steps).",
        "Single CPU thread; no GPU.",
        "Hybrid backend pays for evaluating both networks + disagreement on every step.",
        "85x is conservative — compares against full HTTP–Docker deployment, not bare FMU.",
    ]
    add_bullets(s, items, left=0.8, top=4.5, width=11.7, height=1.6, font_size=15)
    add_takeaway_box(s,
        "8 hours of live BOPTEST training fits in ≈ 5 minutes on the hybrid backend.")
    add_footer(s, idx, SLIDES_TOTAL, "Methodology")


def slide_block_divider(prs, idx, label, sublabel):
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                             prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()
    add_text(s, label, left=0.8, top=2.7, width=12, height=1.4,
             font_size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.8), Inches(4.1),
                               Inches(3), Inches(0.07))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    add_text(s, sublabel, left=0.8, top=4.3, width=12, height=1.0,
             font_size=22, italic=True, color=ACCENT, align=PP_ALIGN.LEFT)


# --- Block 1 slides ---

def slide_block1_v3_v35(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 1 — v3 vs v3.5",
                  subtitle="Two surrogates, two roles")
    headers = ["", "v3", "v3.5"]
    rows = [
        ("Type", "Black-box FFN", "RC-NeuralODE grey-box"),
        ("Explicit C_zon", "No", "Yes (4.413e+05 J/K)"),
        ("Training corpus", "v3 hourly direct-TSup", "v3.5 prepared 15-min"),
        ("Primary strength", "Smooth dynamics; great for RL", "Predictive validity; physical interpretation"),
        ("Primary weakness", "No physical interpretation", "Fails as standalone RL environment"),
    ]
    add_table(s, headers, rows, left=0.8, top=1.6, width=11.7, height=3.3)
    items = [
        "v3 is comfort-oriented; trained to predict T_zone_next and P_total from direct-T_supply data.",
        "v3.5 introduces a learnable C_zon and an explicit dT/dt = (q_net − Q_wall) / C_zon equation.",
        "The two are not competitors — they will be combined in Block 2 as a hybrid.",
    ]
    add_bullets(s, items, left=0.8, top=5.1, width=11.7, height=1.4, font_size=14)
    add_footer(s, idx, SLIDES_TOTAL, "Block 1")


def slide_stage_abc(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Stage A / B / C inverse calibration",
                  subtitle="Hou-and-Evins-compliant three-stage pipeline")
    stages = [
        ("Stage A — Telemetry preprocessing",
         "Latency compensation • Temperature bias removal • Power affine normalization • Rolling denoise • Causal Δt recomputation"),
        ("Stage B — C_zon identification",
         "Excitation-window subselection (top quantile on |ΔT|) • Episode-aware 80/20 split • Optimization with frozen residual heads"),
        ("Stage C — Residual head calibration",
         "C_zon frozen at the Stage B solution • Heads (q_net, p_total) refined • Final RMSE ≈ 0.232 °C (bestest_air)"),
    ]
    y = 1.5
    for name, desc in stages:
        # box
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), Inches(y),
                                  Inches(11.7), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = BLUE
        box.line.width = Pt(1.5)
        add_text(s, name, left=1.0, top=y + 0.12, width=11.3, height=0.45,
                 font_size=17, bold=True, color=BLUE)
        add_text(s, desc, left=1.0, top=y + 0.55, width=11.3, height=0.9,
                 font_size=13, color=INK)
        y += 1.7
    add_footer(s, idx, SLIDES_TOTAL, "Block 1")


def slide_replicative_validity(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 1 — Replicative validity",
                  subtitle="One-step accuracy on the held-out prepared corpus")
    add_image_centered(s, FIG_REAL / "block1_replicative_validity_bars.png",
                       top=1.45, max_height=4.7)
    add_takeaway_box(s,
        "Calibration cuts temperature RMSE by 38% (0.374 → 0.232 °C) and power MAE by 40% (808 → 482 W).")
    add_footer(s, idx, SLIDES_TOTAL, "Block 1")


def slide_predictive_validity(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 1 — Predictive validity across horizons",
                  subtitle="1 h, 4 h, 8 h, 24 h rollout RMSE on held-out episodes")
    add_image_centered(s, FIG_REAL / "block1_predictive_validity_horizon_lines.png",
                       top=1.45, max_height=4.7)
    add_takeaway_box(s,
        "Calibrated RMSE is nearly flat in horizon (0.655 → 0.644 °C from 1 h to 24 h) — the identified C_zon does not accumulate error.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 1")


def slide_rollout_trace(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 1 — 24-hour rollout trace",
                  subtitle="Ground truth (BOPTEST) vs raw v3.5 vs calibrated v3.5")
    add_image_centered(s, FIG_REAL / "block1_rollout_24h_temperature_trace.png",
                       top=1.45, max_height=5.0)
    add_footer(s, idx, SLIDES_TOTAL, "Block 1")


def slide_residual_histograms(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 1 — Residual analysis",
                  subtitle="Stage A/B/C centres the residual distribution and narrows the spread")
    add_image_centered(s, FIG_REAL / "block1_temperature_residual_histograms.png",
                       top=1.45, max_height=4.7)
    add_takeaway_box(s,
        "Residual mean shifts to ≈ 0 °C; σ narrows; p95(|residual|) drops by ~50%.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 1")


def slide_calibration_paradox(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "The calibration paradox",
                  subtitle="Best predictive surrogate ⇏ best RL training environment")
    add_image_centered(s, FIG_REAL / "main_fig3_fidelity_to_rl_gap.png",
                       top=1.45, max_height=4.7)
    add_takeaway_box(s,
        "Calibrated v3.5 wins predictive validity (0.64 °C @ 24h) but loses live closed-loop transfer (~4.4 °C) — a ≈ 7× gap.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 1")


def slide_block1_summary(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 1 — Summary",
                  subtitle="What is established and what is open at the end of Block 1")
    items_left = [
        "Calibrated v3.5 is a strong predictive twin.",
        "C_zon = 4.413e+05 J/K is the identified bestest_air thermal mass.",
        "Stage A/B/C is Hou-and-Evins L3-compliant.",
        "Predictive RMSE flat across 1–24h rollouts.",
    ]
    items_right = [
        "Calibrated v3.5 fails as a stand-alone RL training environment.",
        "Live closed-loop transfer RMSE > 4 °C.",
        "Failure is structural, not a calibration defect.",
        "Block 2 must resolve this without giving up the physics.",
    ]
    add_text(s, "Established", left=0.6, top=1.5, width=6, height=0.4,
             font_size=18, bold=True, color=SUCCESS)
    add_bullets(s, items_left, left=0.6, top=2.0, width=6.0, height=3.5, font_size=15)
    add_text(s, "Open / motivating", left=6.8, top=1.5, width=6, height=0.4,
             font_size=18, bold=True, color=WARN)
    add_bullets(s, items_right, left=6.8, top=2.0, width=6.0, height=3.5, font_size=15)
    add_takeaway_box(s,
        "Block 1 isolates the fidelity-to-RL gap as a boundary condition. Block 2 builds on top of this.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 1")


# --- Block 2 slides ---

def slide_hybrid_resolves(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — Hybrid construction resolves the gap",
                  subtitle="Live BOPTEST transfer RMSE on peak / typical heat windows")
    rows = [
        ("v3 (control-oriented)",       "0.894", "0.745", "good transfer, weak physics"),
        ("v3.5 calibrated (alone)",     "4.320", "4.401", "FAILS as RL environment"),
        ("hybrid_l010 (canonical)",     "0.633", "0.612", "best of both worlds"),
    ]
    add_table(s,
              headers=["Backend", "Peak RMSE (°C)", "Typical RMSE (°C)", "Note"],
              rows=rows,
              left=0.8, top=1.7, width=11.7, height=2.5)
    add_graph_note(s,
        "How to read: lower RMSE is better. The calibrated twin is accurate open-loop but fails closed-loop; the hybrid restores live control accuracy.",
        top=4.15)
    items = [
        "Hybrid keeps the smooth v3 dynamics and the physical anchor of v3.5 simultaneously.",
        "λ_temp = 0.10 and λ_power = 5e-5 are the only two hybrid hyperparameters.",
        "First-divergence step on the typical window: v3 = 1, v3.5 = 1, hybrid = 16.",
    ]
    add_bullets(s, items, left=0.8, top=4.65, width=11.7, height=1.35, font_size=14)
    add_takeaway_box(s,
        "Hybrid recovers the RL-environment property without losing the physics — and is 7× better than v3.5 on transfer.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_warmstart(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — Negative control",
                  subtitle="Direct v3.5 warm-start is worse than scratch")
    add_graph_note(s,
        "How to read: left trace shows comfort behavior; right bars compare m_s/energy. Warm-start is a negative control, not a success case.",
        top=1.28)
    add_image_centered(s, FIG_REAL / "block2_warmstart_negative_eval_kpis.png",
                       top=1.82, max_height=4.35)
    add_takeaway_box(s,
        "Pretraining on calibrated v3.5 then fine-tuning on BOPTEST gives WORSE m_s than scratch. The hybrid path is necessary, not optional.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_thermostatic(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — Thermostatic hybrid (λ_temp = 0.10)",
                  subtitle="Pure v3 vs hybrid_l010 on peak / typical heat windows")
    add_graph_note(s,
        "How to read: compare pure v3 against hybrid_l010 across m_s, violation, RMSE, and energy. Hybrid is strongest in the typical window.",
        top=1.28)
    add_image_centered(s, FIG_REAL / "block2_thermostatic_pure_v3_vs_hybrid_kpis.png",
                       top=1.82, max_height=4.35)
    add_takeaway_box(s,
        "Hybrid wins typical-window m_s (0.041 vs 0.095) and energy. Competitive on peak. Canonical thermostatic point.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_hdrl_sweep(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — HDRL λ_temp sweep",
                  subtitle="Hierarchical controllers REJECT the temperature anchor")
    add_graph_note(s,
        "How to read: each point is one lambda_temp. m_s and violation increase as lambda_temp rises; best HDRL is lambda_temp=0.00.",
        top=1.28)
    add_image_centered(s, FIG_REAL / "block2_hdrl_lambda_sweep_sensitivity.png",
                       top=1.82, max_height=4.35)
    add_takeaway_box(s,
        "Monotone degradation as λ_temp rises. Best HDRL = λ_temp = 0. The high/low-level conflict makes a shared temperature target harmful.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_hdrl_tracking(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — HDRL tracking on winter peak",
                  subtitle="High-level setpoint vs zone temperature (λ_temp = 0)")
    add_graph_note(s,
        "How to read: black is zone temperature; dashed line is HDRL high-level setpoint; shaded band is comfort. This is the best HDRL setting.",
        top=1.28)
    add_image_centered(s, FIG_REAL / "block2_hdrl_l000_winter_tracking.png",
                       top=1.82, max_height=4.65)
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_morl_5d_17d(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — MORL 5-D vs 17-D",
                  subtitle="Same backend, same hybrid loss; observation interface alone changes the outcome")
    add_graph_note(s,
        "How to read: the radar compares normalized KPI axes. The 5-D observation fails; the 17-D TSup-style interface restores viable control.",
        top=1.28)
    add_image_centered(s, FIG_REAL / "block2_morl_5d_vs_17d_radar.png",
                       top=1.82, max_height=4.35)
    add_takeaway_box(s,
        "5-D MORL: m_s = 1.05, violation 75%, RMSE 4.96 °C.  17-D MORL: m_s = 0.099, violation 4.9%, RMSE 0.72 °C.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_morl_pareto(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — MORL Pareto front",
                  subtitle="5 preference vectors + BOPTEST built-in PI reference")
    add_graph_note(s,
        "How to read: x-axis is yearly energy, y-axis is m_s; lower m_s is safer. Preference weights move the policy along the trade-off curve.",
        top=1.28)
    add_image_centered(s, FIG_REAL / "block2_morl_pareto_energy_vs_ms.png",
                       top=1.82, max_height=4.55)
    add_takeaway_box(s,
        "Comfort-leaning preferences satisfy a 5% deployment threshold; the (0, 1) endpoint is a designed safety collapse.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_morl_n5_falsification(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — MORL canonical N=5 seed analysis",
                  subtitle="Pre-registered falsification test")
    rows = [
        ("Neutral canonical w=(0.50, 0.50, 0.00)",  "0.187 ± 0.078", "0.418", "high-variance regime"),
        ("Practical canonical w=(0.75, 0.25, 0.00)", "0.139 ± 0.085", "0.613", "high-variance regime"),
    ]
    add_table(s,
              headers=["Canonical preference", "m_s (mean ± σ, N=5)", "σ / μ", "Status"],
              rows=rows,
              left=0.8, top=1.7, width=11.7, height=1.7)
    items = [
        "Pre-registered action-saturation hypothesis: σ(m_s | Feb_Winter, N=5) < 0.005.",
        "Observed at N=5: σ(m_s | Feb_Winter) = 0.168.  Hypothesis FALSIFIED.",
        "BOPTEST replay test: bit-identical determinism for a fixed checkpoint.",
        "Observed variance is therefore actual policy variance, not simulator noise.",
        "Pre-registered hypothesis and falsification are committed in the audit chain (commit 62dc859d).",
    ]
    add_bullets(s, items, left=0.8, top=3.7, width=11.7, height=2.5, font_size=14)
    add_takeaway_box(s,
        "Falsified hypothesis reported transparently. The N=3-only seasonal-inversion pattern did not survive at N=5.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_seasonal_inversion(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — Seasonal variance heatmap",
                  subtitle="Neutral vs practical canonical at month-level")
    add_graph_note(s,
        "How to read: cells are monthly seed standard deviations of m_s. The initial N=3 seasonal mechanism did not survive N=5 testing.",
        top=1.28)
    add_image_centered(s, FIG_REAL / "block2_morl_seasonal_variance_inversion.png",
                       top=1.82, max_height=4.55)
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_morl_yearly(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — MORL 17-D yearly heatmap",
                  subtitle="Canonical operating point across 12 months")
    add_graph_note(s,
        "How to read: rows are KPI components by month. The heatmap shows where the canonical MORL policy succeeds and where risk remains.",
        top=1.28)
    add_image_centered(s, FIG_REAL / "block2_morl_17d_seasonal_heatmap.png",
                       top=1.82, max_height=4.55)
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_controller_family_lambda(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 2 — Controller-family-specific λ",
                  subtitle="One λ does NOT fit all")
    rows = [
        ("Thermostatic PPO", "0.10", "Low-dim flat observation needs an external anchor"),
        ("HDRL",             "0.00", "Inter-layer conflict; T-anchor harms hierarchical decomposition"),
        ("MORL 17-D",        "0.00", "Rich observation acts as its own self-regularizer"),
    ]
    add_table(s,
              headers=["Controller family", "Optimal λ_temp", "Mechanism"],
              rows=rows,
              left=0.8, top=1.7, width=11.7, height=2.0)
    items = [
        "The optimal λ depends on observation geometry and on the controller's internal hierarchy.",
        "Choosing λ a priori is therefore not a universal hyperparameter — it is a structural property of the controller.",
        "Reported as the central cross-controller finding of Block 2.",
    ]
    add_bullets(s, items, left=0.8, top=4.2, width=11.7, height=1.7, font_size=15)
    add_footer(s, idx, SLIDES_TOTAL, "Block 2")


def slide_negative_results(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Negative results are part of the contribution",
                  subtitle="Failed hypotheses were kept in the final story, not hidden")
    rows = [
        ("Direct v3.5 as RL environment",
         "FAILED", "Live closed-loop RMSE above 4 C despite strong 24h predictive RMSE."),
        ("v3.5 warm-start",
         "FAILED", "Warm-start from calibrated v3.5 is worse than scratch BOPTEST fine-tune."),
        ("HDRL temperature regularization",
         "FAILED", "lambda_temp > 0 monotonically degrades HDRL; best value is 0.00."),
        ("MORL 5-D observation path",
         "FAILED", "m_s=1.046, violation=74.5%, RMSE=4.96 C; 17-D path fixes the interface."),
        ("MORL action-saturation hypothesis",
         "FALSIFIED", "Pre-registered Feb_Winter sigma <0.005 prediction became 0.168 at N=5."),
        ("Block 3 frozen-controller transfer",
         "NOT DEPLOYMENT-READY", "Surrogate transfers across hydronics; frozen controller has comfort/energy failure modes."),
    ]
    add_table(s, headers=["Claim / experiment", "Outcome", "Evidence"],
              rows=rows, left=0.45, top=1.35, width=12.45, height=4.75)
    add_takeaway_box(s,
        "The scientific claim is stronger because the protocol records where the method breaks: fidelity, control utility, and transferability are different properties.")
    add_footer(s, idx, SLIDES_TOTAL, "Negative results")


# --- Block 3 slides ---

def slide_block3_preregistration(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 3 — Pre-registered transferability protocol",
                  subtitle="Manifest committed BEFORE any non-bestest_air BOPTEST run")
    items = [
        "Three testcase candidates pre-selected against four explicit criteria.",
        "Three recalibration regimes pre-defined: none / partial (Stage C only) / full (Stage A/B/C).",
        "Single passfail criterion: m_s_RL ≤ 1.25 × m_s_PI on testcase-specific yearly evaluation.",
        "Four falsifiable hypotheses logged with a priori probability estimates.",
        "Bounded extension policy: N=3 max per cell; no N=5 cascade inside Block 3.",
        "Early-termination clause if H3_weak falsifies on the easiest testcase.",
        "Audit anchor: commit SHA of first manifest commit; locked into configs/block3_testcase_manifest.yaml.",
    ]
    add_bullets(s, items, left=0.8, top=1.5, width=11.7, height=4.5, font_size=15)
    add_takeaway_box(s,
        "Pre-registered design means any outcome (PASS / FAIL / partial) is publishable. We do not depend on a favourable result.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 3")


def slide_block3_testcases(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 3 — Three testcases × three regimes",
                  subtitle="Hydronic family chosen for structural distance from bestest_air")
    rows = [
        ("bestest_hydronic_heat_pump",
         "Primary",  "Same envelope; hydronic loop driven by heat pump"),
        ("bestest_hydronic",
         "Secondary","Same envelope; boiler + radiators"),
        ("singlezone_commercial_hydronic",
         "Stretch",  "Substantially larger envelope; hydronic distribution"),
    ]
    add_table(s,
              headers=["Testcase", "Role", "Structural difference from bestest_air"],
              rows=rows,
              left=0.8, top=1.5, width=11.7, height=2.0)
    rows2 = [
        ("none",    "Frozen recipe deployed via adapter",      "~30 min"),
        ("partial", "Stage C only; C_zon frozen",              "~90 min"),
        ("full",    "Stage A + B + C from scratch on target",  "~240 min"),
    ]
    add_table(s,
              headers=["Regime", "What is recalibrated", "Compute / testcase"],
              rows=rows2,
              left=0.8, top=4.0, width=11.7, height=1.8)
    add_takeaway_box(s,
        "The controller is frozen in ALL three regimes; only the surrogate side is recalibrated.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 3")


def slide_block3_per_testcase(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 3 — Per-testcase results",
                  subtitle="Controller verdict (frozen) and surrogate-side diagnostic (full Stage A/B/C)")
    rows = [
        ("bestest_hydronic_heat_pump",
         "FAIL (m_s 0.665 > 0.579)", "60.2%", "1.89x"),
        ("bestest_hydronic",
         "FAIL (m_s 0.976 > 0.938)", "87.4%", "1.95x"),
        ("singlezone_commercial_hydronic",
         "THRESHOLD PASS (+35.3% energy)", "87.8%", "1.91x"),
    ]
    add_table(s,
              headers=["Testcase", "Controller mode=none verdict",
                       "Surrogate RMSE_T improvement (full)", "C_zon ratio vs bestest_air"],
              rows=rows,
              left=0.5, top=1.7, width=12.3, height=2.5)
    items = [
        "Surrogate side: uniformly transferable with full Stage A/B/C across all three testcases.",
        "Controller side: regime-dependent failure — comfort violation on fast dynamics, energy inflation on slow dynamics.",
        "C_zon consistency: 1.91 ± 0.03 across the hydronic family, independent of actuator type and building scale.",
    ]
    add_bullets(s, items, left=0.8, top=4.6, width=11.7, height=1.7, font_size=14)
    add_footer(s, idx, SLIDES_TOTAL, "Block 3")


def slide_block3_verdict_heatmap(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 3 — Transfer verdict matrix",
                  subtitle="Controller verdict per (testcase, regime)")
    add_image_centered(s, FIG_REAL / "main_fig5_block3_transfer_verdict_heatmap.png",
                       top=1.45, max_height=4.9)
    add_footer(s, idx, SLIDES_TOTAL, "Block 3")


def slide_block3_czon(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 3 — C_zon consistency across testcases",
                  subtitle="Inverse calibration recovers a stable physical parameter")
    add_image_centered(s, FIG_REAL / "main_fig6_block3_czon_consistency.png",
                       top=1.45, max_height=4.9)
    add_takeaway_box(s,
        "C_zon ratio is 1.89×, 1.95×, 1.91× — spread of 3.2%. The Hou-Evins pipeline is robust under cross-testcase re-identification.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 3")


def slide_block3_transfer_gap(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 3 — Hybrid transfer gap on bestest_air vs hydronic",
                  subtitle="The frozen controller does not understand non-T_supply actuators")
    add_image_centered(s, FIG_AUX / "hybrid_transfer_gap_comparison.png",
                       top=1.45, max_height=4.7)
    add_takeaway_box(s,
        "The transferability boundary is the controller-adapter interface, NOT the surrogate's physics representation.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 3")


def slide_block3_decomposition(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Block 3 — Component decomposition finding",
                  subtitle="The hybrid recipe is decomposable; the two components transfer on different terms")
    items_left = [
        "Surrogate component:",
        "  • Transferable on N=3 hydronic testcases under full Stage A/B/C",
        "  • Consistent C_zon identification (1.91 ± 0.03)",
        "  • 60–88% RMSE_T improvement uniformly",
        "  • Hou-Evins pipeline itself is testcase-portable",
    ]
    items_right = [
        "Controller component:",
        "  • Not deployment-ready under frozen-controller scope",
        "  • Regime-dependent failure modes:",
        "    – Fast dynamics → comfort violation",
        "    – Slow dynamics → energy inflation",
        "  • Bottleneck: controller-adapter interface, not physics",
    ]
    add_text(s, "Surrogate-component", left=0.6, top=1.5, width=6.0, height=0.5,
             font_size=18, bold=True, color=SUCCESS)
    add_bullets(s, items_left[1:], left=0.6, top=2.0, width=6.0, height=3.0, font_size=14)
    add_text(s, "Controller-component", left=6.8, top=1.5, width=6.0, height=0.5,
             font_size=18, bold=True, color=WARN)
    add_bullets(s, items_right[1:], left=6.8, top=2.0, width=6.0, height=3.0, font_size=14)
    add_takeaway_box(s,
        "Same nominal hypothesis (H3_weak), opposite outcomes per component. This asymmetry is the central scientific content of Block 3.")
    add_footer(s, idx, SLIDES_TOTAL, "Block 3")


# --- Closing slides ---

def slide_audit_chain(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Pre-registration audit chain",
                  subtitle="Four git-anchored commits preserve the timeline")
    rows = [
        ("1",  "MORL canonical pre-registration",
         "93df9b3", "Hypotheses + stopping rule logged BEFORE seeds 45/46"),
        ("2",  "MORL post-N=5 falsification",
         "62dc859", "Action-saturation hypothesis falsified; reported transparently"),
        ("3",  "Block 3 pre-registration",
         "first manifest commit", "Three testcases, three regimes, four hypotheses logged BEFORE any non-bestest_air run"),
        ("4",  "Block 3 closure",
         "aggregated_results.closed", "N=3 verdicts appended; manifest body bit-identical to first commit"),
    ]
    add_table(s,
              headers=["#", "Anchor", "Commit", "Content"],
              rows=rows,
              left=0.4, top=1.6, width=12.5, height=2.8)
    items = [
        "Every Q1-relevant claim in this work is traceable to one of these four commits.",
        "Server-side GitHub timestamps make the chain externally verifiable.",
    ]
    add_bullets(s, items, left=0.6, top=4.6, width=12.0, height=1.0, font_size=14)
    add_takeaway_box(s,
        "Most DRL-HVAC papers have zero pre-registration anchors. This work has four.")
    add_footer(s, idx, SLIDES_TOTAL, "Audit")


def slide_hou_evins(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Hou-and-Evins compliance",
                  subtitle="11 supplementary tables S1–S11 cover the full protocol")
    rows = [
        ("Stage 1 (Sample generation)", "S1, S2, S3", "L3"),
        ("Stage 2 (Data processing)",   "S4, S5, S6, S7", "L3"),
        ("Stage 3 (Architecture & training)", "S8, S9, S10", "L3"),
        ("Stage 4 (Validity)",          "S11 + Tables 2, 3", "L3"),
    ]
    add_table(s,
              headers=["Hou-and-Evins protocol stage", "Supplementary tables",
                       "Reporting level"],
              rows=rows,
              left=0.8, top=1.7, width=11.7, height=2.4)
    items = [
        "Every methodology choice has a numerical justification table.",
        "15 of 17 protocol requirements satisfied at Reporting Level 3.",
        "Two L2 cells explicitly noted as known limitations (excitation-window prose, scaling justification).",
        "Most DRL-HVAC studies in the literature provide L1 reporting or less.",
    ]
    add_bullets(s, items, left=0.8, top=4.5, width=11.7, height=1.8, font_size=15)
    add_footer(s, idx, SLIDES_TOTAL, "Methodology audit")


def slide_limitations(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Limitations",
                  subtitle="Boundaries of the claim, stated before the conclusion")
    rows = [
        ("Scope", "Single primary weather profile and BOPTEST testcase family; no universal climate or multi-zone claim."),
        ("MORL stability", "Both canonical MORL points remain high-variance at N=5; no deployment-stable MORL claim."),
        ("Transferability", "Block 3 proves surrogate-component transfer, not frozen-controller deployment readiness."),
        ("Metric threshold", "m_s threshold can mask energy penalties: commercial hydronic passes threshold but uses +35.3% energy vs PI."),
        ("Actuator adapter", "Hydronic transfer is adapter-mediated, not literal direct-T_supply transfer."),
        ("Speed benchmark", "85x speed-up is against BOPTEST RTE HTTP-Docker loop, not a bare-FMU microbenchmark."),
    ]
    add_table(s, headers=["Boundary", "What is not claimed"],
              rows=rows, left=0.6, top=1.35, width=12.1, height=4.6)
    add_takeaway_box(s,
        "These limitations are not footnotes: they define the exact scientific scope and the natural Block 4 follow-up.")
    add_footer(s, idx, SLIDES_TOTAL, "Limitations")


def slide_threats(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Threats to validity",
                  subtitle="Explicit limitations, transparently stated")
    items = [
        "Single weather file: bestest_air TMY only. Climate generalization not tested.",
        "Single testcase family for cross-transfer: hydronic, 3 BOPTEST testcases.",
        "MORL controller variance high at N=5 (σ/μ = 0.42–0.61); deployment-stability not yet established.",
        "Action-saturation hypothesis falsified — N=3 seasonal-inversion pattern did not survive at N=5.",
        "Threshold-based PASS/FAIL framework can mask multi-objective trade-offs (commercial PASS with +35% energy).",
        "Speed comparison is against BOPTEST RTE HTTP–Docker deployment, not bare FMU; FMU-direct benchmark deferred to Linux runtime.",
        "Controller fine-tune on target-recalibrated surrogate explicitly scope-excluded — natural Block 4 follow-up.",
    ]
    add_bullets(s, items, left=0.6, top=1.5, width=12.0, height=5.0, font_size=15)
    add_footer(s, idx, SLIDES_TOTAL, "Threats")


def slide_conclusion(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Conclusion — main claim",
                  subtitle="Component-level decomposition of surrogate-based HVAC RL transferability")
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(1.5),
                              Inches(12.3), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = INK
    box.line.width = Pt(1.5)
    quote = (
        "On the BOPTEST hydronic family (N = 3 testcases), the Hou-and-Evins "
        "three-stage inverse calibration pipeline transfers the SURROGATE "
        "component of the hybrid recipe with consistent physical-parameter "
        "identification (C_zon ratio 1.91 ± 0.03). The frozen "
        "direct-T_supply CONTROLLER, routed through a mechanical hydronic "
        "adapter, does not transfer in a deployment-ready sense; its "
        "failure manifests as comfort violation under fast envelope "
        "dynamics and as energy inflation under slow envelope dynamics.\n\n"
        "The transferability bottleneck is the controller-adapter "
        "interface, not the surrogate's physics representation capacity."
    )
    add_text(s, quote, left=0.8, top=1.7, width=11.7, height=3.2,
             font_size=14, italic=True, color=INK, align=PP_ALIGN.LEFT)
    items = [
        "Speed-up: 85x over BOPTEST RTE HTTP–Docker deployment.",
        "Controller-family-specific λ established for thermostatic / HDRL / MORL.",
        "Pre-registered three-tier protocol with four audit anchors in git.",
        "Hou-and-Evins L3 numerical justification (S1–S11).",
    ]
    add_bullets(s, items, left=0.6, top=5.2, width=12.0, height=1.5, font_size=15)
    add_footer(s, idx, SLIDES_TOTAL, "Conclusion")


def slide_future_work(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Future work — Block 4",
                  subtitle="Natural follow-up identified by Block 3 component decomposition")
    items_left = [
        "Block 4: controller fine-tune on target-recalibrated surrogate",
        "  • Surrogate is already transferable (proved in Block 3)",
        "  • Controller is the bottleneck (proved in Block 3)",
        "  • Fine-tune step closes the gap by construction",
    ]
    items_right = [
        "Adjacent directions (separate papers):",
        "  • Multi-zone testcases",
        "  • Cross-climate generalization",
        "  • Continual building learning (LEGION / DPMM)",
        "  • Differentiable world-model RL (Dreamer-style)",
        "  • Hierarchical preference architectures (HERON)",
    ]
    add_text(s, "Pre-identified primary follow-up", left=0.6, top=1.5, width=6.0,
             height=0.5, font_size=18, bold=True, color=SUCCESS)
    add_bullets(s, items_left[1:], left=0.6, top=2.0, width=6.0, height=3.0, font_size=14)
    add_text(s, "Other future directions", left=6.8, top=1.5, width=6.0,
             height=0.5, font_size=18, bold=True, color=BLUE)
    add_bullets(s, items_right[1:], left=6.8, top=2.0, width=6.0, height=3.0, font_size=14)
    add_takeaway_box(s,
        "Block 4 scope is identified, not just future-handwaving — the bottleneck is named and the prerequisite is established.")
    add_footer(s, idx, SLIDES_TOTAL, "Future work")


def slide_acknowledgments(prs, idx):
    s = blank_slide(prs)
    add_title_bar(s, "Thank you",
                  subtitle="Questions, discussion, defense Q&A")
    add_text(s, "Questions?", left=0.8, top=2.0, width=12.0, height=1.5,
             font_size=72, bold=True, color=INK, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    items_resources = [
        "Repository: see roadmap.md in project root",
        "Audit chain: 93df9b3 → 62dc859 → Block 3 manifest → Block 3 closure",
        "Manuscript: docs/hvac_paper_skeleton_q1_restructured.docx",
        "Supplementary tables: reports/hou_evins_*.csv (S1–S11)",
    ]
    add_text(s, "Resources", left=0.6, top=4.5, width=12.0, height=0.4,
             font_size=15, bold=True, color=MUTED)
    add_bullets(s, items_resources, left=0.6, top=4.9, width=12.0, height=2.0,
                font_size=13, color=MUTED)
    add_footer(s, idx, SLIDES_TOTAL, "Q&A")


# ===========================================================================
# Build sequence
# ===========================================================================

def build():
    prs = make_presentation()
    # Concise main-body defense deck.
    #
    # Removed from the 53-slide full version:
    # - duplicate motivation/methodology dividers,
    # - secondary diagnostic plots already covered by stronger aggregate figures,
    # - per-family drill-down slides that can be discussed verbally,
    # - separate threats/future-work/Hou-Evins slides folded into limitations/conclusion.
    #
    # The retained spine keeps every central claim: surrogate fidelity, the
    # fidelity-to-RL gap, hybrid regularization, negative findings, MORL seed
    # variance, Block 3 transferability, audit, and limitations.
    slides = [
        ("title",              slide_title),
        ("outline",            slide_outline),
        ("relevance",          slide_relevance),
        ("research_questions", slide_research_questions),
        ("scientific_novelty", slide_scientific_novelty),
        ("related_works",      lambda p, i: slide_empty_placeholder(p, i,
            "Related Works",
            "[Empty by design.]\n\nAuthor will populate this slide with positioning against prior DRL-HVAC, surrogate-modeling, physics-informed ML, and BOPTEST-benchmarking literature.")),
        ("literature_review",  lambda p, i: slide_empty_placeholder(p, i,
            "Literature Review",
            "[Empty by design.]\n\nAuthor will populate this slide with the systematic literature review and the position-of-this-work paragraph.")),
        ("pipeline_schematic", slide_pipeline_schematic),
        ("hybrid_loss",        slide_hybrid_loss),
        ("surrogate_architecture", slide_surrogate_architecture),
        ("surrogate_equations", slide_surrogate_equations),
        ("control_interface_reward", slide_control_interface_reward),
        ("controller_nn_details", slide_controller_nn_details),
        ("speed_benchmark",    slide_speed_benchmark),
        ("block1_v3_v35",      slide_block1_v3_v35),
        ("stage_abc",          slide_stage_abc),
        ("predictive",         slide_predictive_validity),
        ("rollout_trace",      slide_rollout_trace),
        ("calibration_paradox", slide_calibration_paradox),
        ("hybrid_resolves",    slide_hybrid_resolves),
        ("warmstart",          slide_warmstart),
        ("hdrl_sweep",         slide_hdrl_sweep),
        ("morl_5d_17d",        slide_morl_5d_17d),
        ("morl_pareto",        slide_morl_pareto),
        ("morl_n5",            slide_morl_n5_falsification),
        ("controller_lambda",  slide_controller_family_lambda),
        ("negative_results",   slide_negative_results),
        ("block3_prereg",      slide_block3_preregistration),
        ("block3_testcases",   slide_block3_testcases),
        ("block3_per_testcase", slide_block3_per_testcase),
        ("block3_heatmap",     slide_block3_verdict_heatmap),
        ("block3_czon",        slide_block3_czon),
        ("block3_decomposition", slide_block3_decomposition),
        ("audit_chain",        slide_audit_chain),
        ("limitations",        slide_limitations),
        ("conclusion",         slide_conclusion),
        ("acknowledgments",    slide_acknowledgments),
    ]
    global SLIDES_TOTAL
    SLIDES_TOTAL = len(slides)
    for i, (_, builder) in enumerate(slides, start=1):
        builder(prs, i)

    prs.save(OUT_PATH)
    print(f"[OK] {OUT_PATH} ({SLIDES_TOTAL} slides)")


if __name__ == "__main__":
    build()
